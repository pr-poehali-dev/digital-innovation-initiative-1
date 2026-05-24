"""
Загрузка и управление документами (PDF, DOCX, PPTX).
Извлечение текста и структуры из файлов.
"""
import json
import os
import base64
import io
import psycopg2
import boto3


def get_db():
    conn = psycopg2.connect(os.environ["DATABASE_URL"])
    conn.autocommit = False
    return conn


def get_schema():
    return os.environ.get("MAIN_DB_SCHEMA", "public")


def cors_headers():
    return {
        "Access-Control-Allow-Origin": "*",
        "Access-Control-Allow-Methods": "GET, POST, PUT, DELETE, OPTIONS",
        "Access-Control-Allow-Headers": "Content-Type, X-Session-Id",
    }


def json_response(data, status=200):
    return {
        "statusCode": status,
        "headers": {**cors_headers(), "Content-Type": "application/json"},
        "body": json.dumps(data, ensure_ascii=False, default=str),
    }


def get_current_user(conn, session_id):
    if not session_id:
        return None
    schema = get_schema()
    cur = conn.cursor()
    cur.execute(
        f"SELECT u.id, u.email, u.name FROM {schema}.sessions s JOIN {schema}.users u ON u.id = s.user_id WHERE s.id = %s AND s.expires_at > NOW()",
        (session_id,),
    )
    row = cur.fetchone()
    if row:
        return {"id": row[0], "email": row[1], "name": row[2]}
    return None


def get_s3():
    return boto3.client(
        "s3",
        endpoint_url="https://bucket.poehali.dev",
        aws_access_key_id=os.environ["AWS_ACCESS_KEY_ID"],
        aws_secret_access_key=os.environ["AWS_SECRET_ACCESS_KEY"],
    )


MAX_TEXT_LEN = 500000  # 500 КБ текста — поддержка больших дипломов
CHUNK_SIZE = 1500       # символов на чанк
CHUNK_OVERLAP = 200     # перекрытие чанков для сохранения контекста


def chunk_text(text: str) -> list:
    """Делит длинный текст на перекрывающиеся чанки."""
    chunks = []
    if not text:
        return chunks
    start = 0
    idx = 0
    while start < len(text):
        end = min(start + CHUNK_SIZE, len(text))
        # Стараемся резать по концу предложения
        if end < len(text):
            for sep in [". ", "!\n", "?\n", "\n\n", ".\n", "\n"]:
                cut = text.rfind(sep, start, end)
                if cut > start + CHUNK_SIZE // 2:
                    end = cut + len(sep)
                    break
        chunk = text[start:end].strip()
        if chunk:
            chunks.append({"index": idx, "content": chunk, "page": None})
            idx += 1
        start = end - CHUNK_OVERLAP if end - CHUNK_OVERLAP > start else end
    return chunks


def extract_text_from_pdf(data: bytes):
    """Возвращает (полный текст, чанки с указанием страниц, кол-во страниц)."""
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(data))
        full_text = []
        page_texts = []
        for i, page in enumerate(reader.pages):
            t = page.extract_text() or ""
            page_texts.append((i + 1, t))
            full_text.append(t)
        text = "\n".join(full_text)[:MAX_TEXT_LEN]
        # Чанки с привязкой к страницам
        chunks = []
        idx = 0
        for page_num, page_text in page_texts:
            if not page_text.strip():
                continue
            for c in chunk_text(page_text):
                chunks.append({"index": idx, "content": c["content"], "page": page_num})
                idx += 1
        return text, chunks, len(reader.pages)
    except Exception as e:
        return f"[Ошибка извлечения текста PDF: {e}]", [], 0


def extract_text_from_docx(data: bytes):
    try:
        import docx
        doc = docx.Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs)[:MAX_TEXT_LEN]
        chunks = chunk_text(text)
        return text, chunks, None
    except Exception as e:
        return f"[Ошибка извлечения текста DOCX: {e}]", [], None


def extract_from_pptx_full(data: bytes):
    """Обёртка для PPTX чтобы совпадал интерфейс."""
    result = extract_from_pptx(data)
    chunks = chunk_text(result["text"])
    return result["text"][:MAX_TEXT_LEN], chunks, len(result.get("structure", []))


def extract_from_pptx(data: bytes) -> dict:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides):
            title = ""
            bullets = []
            notes = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.shape_type == 13:
                        continue
                    for j, para in enumerate(shape.text_frame.paragraphs):
                        text = para.text.strip()
                        if not text:
                            continue
                        if j == 0 and not title:
                            title = text
                        else:
                            bullets.append(text)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            slides.append({
                "index": i + 1,
                "title": title,
                "bullets": bullets,
                "notes": notes,
            })
        full_text = "\n\n".join(
            f"Слайд {s['index']}: {s['title']}\n" + "\n".join(s['bullets'])
            for s in slides
        )
        return {"text": full_text[:50000], "structure": slides}
    except Exception as e:
        return {"text": f"[Ошибка извлечения PPTX: {e}]", "structure": []}


def log_activity(cur, schema, project_id, user_id, action, entity_type=None, entity_id=None, details=None):
    cur.execute(
        f"INSERT INTO {schema}.activity_log (project_id, user_id, action, entity_type, entity_id, details) VALUES (%s, %s, %s, %s, %s, %s)",
        (project_id, user_id, action, entity_type, entity_id, details),
    )


def handler(event: dict, context) -> dict:
    if event.get("httpMethod") == "OPTIONS":
        return {"statusCode": 200, "headers": cors_headers(), "body": ""}

    method = event.get("httpMethod", "GET")
    path = event.get("path", "/")
    body = {}
    if event.get("body"):
        try:
            body = json.loads(event["body"])
        except Exception:
            pass

    session_id = event.get("headers", {}).get("X-Session-Id", "")
    conn = get_db()
    schema = get_schema()

    try:
        user = get_current_user(conn, session_id)
        if not user:
            return json_response({"error": "Не авторизован"}, 401)

        cur = conn.cursor()
        path_parts = path.strip("/").split("/")

        # POST /upload — загрузить файл
        if method == "POST" and "upload" in path:
            project_id = body.get("project_id")
            filename = body.get("filename", "file")
            file_data_b64 = body.get("file_data", "")
            file_type = body.get("file_type", "").lower()

            if not project_id or not file_data_b64:
                return json_response({"error": "Не хватает данных"}, 400)

            # Проверка доступа к проекту
            cur.execute(
                f"SELECT role FROM {schema}.project_members WHERE project_id = %s AND user_id = %s",
                (project_id, user["id"]),
            )
            if not cur.fetchone():
                return json_response({"error": "Нет доступа"}, 403)

            file_bytes = base64.b64decode(file_data_b64)
            file_size = len(file_bytes)

            # Загрузить в S3
            s3_key = f"documents/{project_id}/{user['id']}_{filename}"
            s3 = get_s3()
            content_types = {
                "pdf": "application/pdf",
                "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            }
            s3.put_object(
                Bucket="files",
                Key=s3_key,
                Body=file_bytes,
                ContentType=content_types.get(file_type, "application/octet-stream"),
            )

            # Извлечь текст и чанки
            extracted_text = ""
            structure_json = None
            chunks = []
            page_count = None
            category = body.get("category", "other")

            if file_type == "pdf":
                extracted_text, chunks, page_count = extract_text_from_pdf(file_bytes)
            elif file_type == "docx":
                extracted_text, chunks, page_count = extract_text_from_docx(file_bytes)
            elif file_type == "pptx":
                result = extract_from_pptx(file_bytes)
                extracted_text = result["text"][:MAX_TEXT_LEN]
                structure_json = json.dumps(result["structure"], ensure_ascii=False)
                chunks = chunk_text(extracted_text)
                page_count = len(result.get("structure", []))

            cur.execute(
                f"""INSERT INTO {schema}.documents
                    (project_id, uploaded_by, filename, original_name, file_type, file_size, s3_key,
                     extracted_text, structure_json, status, category, page_count, extracted_length)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, 'ready', %s, %s, %s)
                    RETURNING id""",
                (project_id, user["id"], s3_key, filename, file_type, file_size, s3_key,
                 extracted_text, structure_json, category, page_count, len(extracted_text)),
            )
            doc_id = cur.fetchone()[0]

            # Сохраняем чанки
            for ch in chunks[:500]:  # лимит чанков на 1 файл
                cur.execute(
                    f"""INSERT INTO {schema}.document_chunks (document_id, chunk_index, page_number, content, content_length)
                        VALUES (%s, %s, %s, %s, %s)""",
                    (doc_id, ch["index"], ch.get("page"), ch["content"], len(ch["content"])),
                )

            log_activity(cur, schema, project_id, user["id"], "uploaded_document", "document", doc_id, filename)
            conn.commit()

            return json_response({
                "id": doc_id,
                "filename": filename,
                "file_type": file_type,
                "file_size": file_size,
                "status": "ready",
                "category": category,
                "page_count": page_count,
                "chunks_count": len(chunks),
                "text_length": len(extracted_text),
                "has_structure": structure_json is not None,
            })

        # GET /project/{project_id} — список документов проекта
        if method == "GET" and "project" in path:
            for i, part in enumerate(path_parts):
                if part == "project" and i + 1 < len(path_parts):
                    project_id = int(path_parts[i + 1])
                    break
            else:
                return json_response({"error": "Нет project_id"}, 400)

            cur.execute(
                f"SELECT role FROM {schema}.project_members WHERE project_id = %s AND user_id = %s",
                (project_id, user["id"]),
            )
            if not cur.fetchone():
                return json_response({"error": "Нет доступа"}, 403)

            cur.execute(
                f"""SELECT d.id, d.original_name, d.file_type, d.file_size, d.status, d.created_at, u.name,
                    d.category, d.page_count, d.extracted_length
                    FROM {schema}.documents d JOIN {schema}.users u ON u.id = d.uploaded_by
                    WHERE d.project_id = %s ORDER BY d.created_at DESC""",
                (project_id,),
            )
            docs = [
                {
                    "id": r[0], "name": r[1], "file_type": r[2],
                    "file_size": r[3], "status": r[4],
                    "created_at": str(r[5]), "uploaded_by": r[6],
                    "category": r[7] or "other",
                    "page_count": r[8], "text_length": r[9],
                }
                for r in cur.fetchall()
            ]
            return json_response({"documents": docs})

        # PUT /{id}/category — изменить категорию
        if method == "PUT" and "category" in path_parts:
            doc_id = int(path_parts[-2])
            new_cat = (body.get("category") or "other").strip()
            cur.execute(
                f"SELECT project_id FROM {schema}.documents WHERE id = %s",
                (doc_id,),
            )
            row = cur.fetchone()
            if not row:
                return json_response({"error": "Не найдено"}, 404)
            cur.execute(
                f"SELECT role FROM {schema}.project_members WHERE project_id = %s AND user_id = %s",
                (row[0], user["id"]),
            )
            if not cur.fetchone():
                return json_response({"error": "Нет доступа"}, 403)
            cur.execute(
                f"UPDATE {schema}.documents SET category = %s WHERE id = %s",
                (new_cat, doc_id),
            )
            conn.commit()
            return json_response({"ok": True, "category": new_cat})

        # GET /{id}/text — получить извлечённый текст документа
        if method == "GET" and path_parts[-1] == "text":
            doc_id = int(path_parts[-2])
            cur.execute(
                f"SELECT d.extracted_text, d.structure_json, d.project_id FROM {schema}.documents d WHERE d.id = %s",
                (doc_id,),
            )
            row = cur.fetchone()
            if not row:
                return json_response({"error": "Не найдено"}, 404)
            cur.execute(
                f"SELECT role FROM {schema}.project_members WHERE project_id = %s AND user_id = %s",
                (row[2], user["id"]),
            )
            if not cur.fetchone():
                return json_response({"error": "Нет доступа"}, 403)
            structure = None
            if row[1]:
                try:
                    structure = json.loads(row[1])
                except Exception:
                    pass
            return json_response({"text": row[0], "structure": structure})

        return json_response({"error": "Not found"}, 404)

    finally:
        conn.close()