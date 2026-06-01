"""
Educational Passport — паспорт образования пользователя.

Действия (action namespace v1):
  - education.list                — список своих записей (с фильтрами)
  - education.get                 — детали записи
  - education.create              — создать запись вручную
  - education.update              — обновить
  - education.archive             — soft delete
  - education.upload_file         — прикрепить файл (base64) + автозапуск анализа
  - education.analyze             — повторный AI-анализ файла
  - education.confirm             — подтвердить AI-извлечение пользователем
  - education.profile_summary     — счётчики для Dashboard
"""
import json
import os
import uuid
import base64
import logging
from datetime import datetime
import psycopg2

logging.basicConfig(level=logging.INFO)
log = logging.getLogger("education")


ALLOWED_ACTIONS = {
    "education.list",
    "education.get",
    "education.create",
    "education.update",
    "education.archive",
    "education.upload_file",
    "education.get_upload_url",
    "education.file_ready",
    "education.get_file_url",
    "education.analyze",
    "education.confirm",
    "education.profile_summary",
}

KIND_GROUPS = {
    "formal": {"degree", "certificate", "course", "program"},
    "material": {"book", "lecture", "presentation", "methodology", "notes", "article", "material"},
}


INDEXER_URL = os.environ.get("SEARCH_INDEXER_URL", "")
INDEXER_TOKEN = os.environ.get("SEARCH_INDEXER_TOKEN", "")


def get_db():
    conn = psycopg2.connect(os.environ["DATABASE_URL"])
    conn.autocommit = False
    return conn


def notify_indexer(action: str, entity_id: int = None):
    if not INDEXER_URL:
        return
    try:
        import urllib.request
        body = {"entity_type": "education"}
        if entity_id:
            body["entity_id"] = entity_id
        hdrs = {"Content-Type": "application/json"}
        if INDEXER_TOKEN:
            hdrs["X-Internal-Token"] = INDEXER_TOKEN
        req = urllib.request.Request(
            f"{INDEXER_URL}?action={action}",
            data=json.dumps(body).encode(),
            headers=hdrs,
            method="POST",
        )
        urllib.request.urlopen(req, timeout=3)
    except Exception:
        pass


def get_schema():
    return os.environ.get("MAIN_DB_SCHEMA", "public")


ALLOWED_ORIGINS = {
    "https://raven.moscow",
    "https://www.raven.moscow",
    "https://docmind.ai",
    "https://digital-innovation-initiative-1--preview.poehali.dev",
    "https://poehali.dev",
    "http://localhost:5173",
    "http://localhost:3000",
}


def _is_allowed_origin(origin: str) -> bool:
    if not origin:
        return False
    if origin in ALLOWED_ORIGINS:
        return True
    try:
        from urllib.parse import urlparse
        parsed = urlparse(origin)
        if parsed.scheme not in ("https", "http"):
            return False
        hostname = (parsed.hostname or "").lower()
        return hostname == "poehali.dev" or hostname.endswith(".poehali.dev")
    except Exception:
        return False


def cors_headers(origin: str = None):
    """Strict CORS deny-by-default."""
    headers = {
        "Access-Control-Allow-Methods": "POST, OPTIONS",
        "Access-Control-Allow-Headers": "Content-Type, X-Session-Id",
        "Vary": "Origin",
    }
    if _is_allowed_origin(origin):
        headers["Access-Control-Allow-Origin"] = origin
    return headers


def ok_response(data, request_id, origin=None):
    return {
        "statusCode": 200,
        "headers": {**cors_headers(origin), "Content-Type": "application/json", "X-Request-Id": request_id, "X-Api-Version": "v1"},
        "body": json.dumps({"ok": True, "request_id": request_id, "data": data}, ensure_ascii=False, default=str),
    }


def err_response(code, message, status, request_id, origin=None):
    return {
        "statusCode": status,
        "headers": {**cors_headers(origin), "Content-Type": "application/json", "X-Request-Id": request_id, "X-Api-Version": "v1"},
        "body": json.dumps({"ok": False, "request_id": request_id, "error": {"code": code, "message": message}}, ensure_ascii=False),
    }


def get_current_user(conn, session_id):
    if not session_id:
        return None
    schema = get_schema()
    cur = conn.cursor()
    cur.execute(
        f"SELECT u.id FROM {schema}.sessions s JOIN {schema}.users u ON u.id = s.user_id WHERE s.id = %s AND s.expires_at > NOW()",
        (session_id,),
    )
    row = cur.fetchone()
    return {"id": row[0]} if row else None


def get_s3():
    import boto3
    return boto3.client(
        "s3",
        endpoint_url="https://bucket.poehali.dev",
        aws_access_key_id=os.environ["AWS_ACCESS_KEY_ID"],
        aws_secret_access_key=os.environ["AWS_SECRET_ACCESS_KEY"],
    )


# ============================================================
# AI-экстрактор метаданных
# ============================================================


def call_yandex_gpt(messages):
    api_key = os.environ.get("YANDEX_GPT_API_KEY", "")
    folder_id = os.environ.get("YANDEX_FOLDER_ID", "")
    if not api_key or not folder_id:
        return "[AI недоступен: добавьте YANDEX_GPT_API_KEY и YANDEX_FOLDER_ID]"
    import urllib.request, urllib.error
    yandex_messages = [{"role": m["role"], "text": m["content"]} for m in messages]
    payload = json.dumps({
        "modelUri": f"gpt://{folder_id}/yandexgpt/latest",
        "completionOptions": {"stream": False, "temperature": 0.2, "maxTokens": 2000},
        "messages": yandex_messages,
    }).encode()
    req = urllib.request.Request(
        "https://llm.api.cloud.yandex.net/foundationModels/v1/completion",
        data=payload,
        headers={"Authorization": f"Api-Key {api_key}", "Content-Type": "application/json"},
    )
    try:
        with urllib.request.urlopen(req, timeout=90) as resp:
            result = json.loads(resp.read())
            return result["result"]["alternatives"][0]["message"]["text"]
    except Exception as e:
        return f"[Ошибка AI: {e}]"


def extract_text_from_file(file_bytes: bytes, mime: str) -> str:
    """Извлечение текста из PDF/DOCX/PPTX/TXT. Для картинок — пока возвращаем пусто."""
    import io
    try:
        if "pdf" in mime:
            import PyPDF2
            reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            return "\n".join((p.extract_text() or "") for p in reader.pages)[:30000]
        if "wordprocessingml" in mime or mime.endswith("/docx"):
            import docx
            d = docx.Document(io.BytesIO(file_bytes))
            return "\n".join(p.text for p in d.paragraphs if p.text.strip())[:30000]
        if "presentationml" in mime or mime.endswith("/pptx"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            return "\n".join(parts)[:30000]
        if mime.startswith("text/"):
            return file_bytes.decode("utf-8", errors="ignore")[:30000]
        # JPG / PNG / HEIC — скан документа, используем Yandex Vision OCR
        if mime.startswith("image/") or any(mime.endswith(ext) for ext in ("/jpeg", "/jpg", "/png", "/webp", "/heic")):
            return ocr_image_bytes(file_bytes)
        return ""
    except Exception as e:
        return f"[Ошибка извлечения: {e}]"


def ocr_image_bytes(image_bytes: bytes) -> str:
    """OCR через Yandex Vision API для сканов дипломов и сертификатов."""
    api_key = os.environ.get("YANDEX_GPT_API_KEY", "")
    folder_id = os.environ.get("YANDEX_FOLDER_ID", "")
    if not api_key or not folder_id:
        return "[OCR недоступен: нет ключей Yandex]"
    import urllib.request, urllib.error
    b64 = base64.b64encode(image_bytes).decode("ascii")
    payload = json.dumps({
        "folderId": folder_id,
        "analyze_specs": [{
            "content": b64,
            "features": [{"type": "TEXT_DETECTION", "text_detection_config": {"language_codes": ["ru", "en"]}}],
        }],
    }).encode()
    req = urllib.request.Request(
        "https://vision.api.cloud.yandex.net/vision/v1/batchAnalyze",
        data=payload,
        headers={"Authorization": f"Api-Key {api_key}", "Content-Type": "application/json"},
    )
    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            result = json.loads(resp.read())
        lines = []
        for r in result.get("results", []):
            for sub in r.get("results", []):
                td = sub.get("textDetection", {})
                for page in td.get("pages", []):
                    for block in page.get("blocks", []):
                        for line in block.get("lines", []):
                            words = [w.get("text", "") for w in line.get("words", [])]
                            text = " ".join(words).strip()
                            if text:
                                lines.append(text)
        return "\n".join(lines)[:30000] if lines else "[OCR: текст не найден]"
    except Exception as e:
        return f"[Ошибка OCR: {e}]"


def ai_extract_formal(text: str) -> dict:
    """AI извлекает метаданные диплома/сертификата/курса.
    Возвращает dict готовый для записи в БД (extracted_json + поля)."""
    prompt = f"""Перед тобой текст документа об образовании (диплом / сертификат / курс).
Извлеки структурированные метаданные. Отвечай СТРОГО в формате JSON, без markdown-блоков и комментариев.

Текст документа:
\"\"\"
{text[:8000]}
\"\"\"

Формат ответа (валидный JSON):
{{
  "title": "название документа или программы",
  "institution_name": "учреждение / платформа выдачи",
  "field_of_study": "направление / специальность",
  "level": "bachelor | master | phd | professional | online | school | other",
  "issued_at": "YYYY-MM-DD или null",
  "hours": число академических часов или null,
  "grade": "оценка/результат или null",
  "topics": ["изученные дисциплины / темы"],
  "suggested_competencies": ["вероятные компетенции которые получены"],
  "summary": "краткое описание в 1-2 предложения"
}}

ВАЖНО:
- Если поле не нашёл — ставь null или пустой массив
- НЕ фантазируй компетенции которых нет в тексте
- topics и suggested_competencies — максимум 10 элементов
- ответ — ТОЛЬКО валидный JSON, ничего больше"""

    answer = call_yandex_gpt([
        {"role": "system", "content": "Ты — экстрактор метаданных образовательных документов. Возвращаешь только валидный JSON."},
        {"role": "user", "content": prompt},
    ])

    # Парсим JSON из ответа
    import re
    answer = answer.strip()
    # Убираем возможные markdown-обёртки
    answer = re.sub(r"^```(?:json)?\s*", "", answer)
    answer = re.sub(r"\s*```$", "", answer)
    try:
        return json.loads(answer)
    except Exception:
        # Попытка вытащить JSON из текста
        match = re.search(r"\{.*\}", answer, re.DOTALL)
        if match:
            try:
                return json.loads(match.group(0))
            except Exception:
                pass
        return {"error": "Не удалось распарсить ответ AI", "raw": answer[:500]}


def ai_extract_material(text: str) -> dict:
    """AI извлекает метаданные учебного материала (книга/лекция/конспект)."""
    prompt = f"""Перед тобой текст учебного материала (книга, лекция, конспект, методичка, статья).
Извлеки структурированные метаданные. Отвечай СТРОГО валидным JSON.

Текст материала (фрагмент):
\"\"\"
{text[:8000]}
\"\"\"

Формат:
{{
  "title": "название материала",
  "material_type": "book | lecture | presentation | methodology | notes | article | other",
  "author": "автор или null",
  "main_topics": ["основные темы материала"],
  "domain": "область знаний (например: 'юриспруденция', 'управление проектами', 'программирование')",
  "suggested_competencies": ["компетенции которые материал ПОКРЫВАЕТ (не значит что усвоены!)"],
  "summary": "о чём этот материал в 1-2 предложениях"
}}

ОЧЕНЬ ВАЖНО:
- Не пиши «пользователь владеет темой» — мы не знаем уровень освоения
- Можно писать только «материал покрывает темы X, Y, Z»
- main_topics — максимум 10
- ответ — ТОЛЬКО валидный JSON"""

    answer = call_yandex_gpt([
        {"role": "system", "content": "Ты — экстрактор метаданных учебных материалов. Не утверждаешь освоение, только покрытие тем. Возвращаешь только JSON."},
        {"role": "user", "content": prompt},
    ])

    import re
    answer = answer.strip()
    answer = re.sub(r"^```(?:json)?\s*", "", answer)
    answer = re.sub(r"\s*```$", "", answer)
    try:
        return json.loads(answer)
    except Exception:
        match = re.search(r"\{.*\}", answer, re.DOTALL)
        if match:
            try:
                return json.loads(match.group(0))
            except Exception:
                pass
        return {"error": "Не удалось распарсить ответ AI", "raw": answer[:500]}


# ============================================================
# Handlers
# ============================================================


def handle_list(conn, user, body, request_id, origin):
    schema = get_schema()
    cur = conn.cursor()
    kind_filter = body.get("kind_filter")  # all / formal / material / specific kind
    status_filter = body.get("status_filter")  # all / confirmed / needs_review / draft

    sql = f"""SELECT id, kind, title, issuer_name, institution_name, field_of_study,
                level, issued_at, status, study_status, source_type, is_confirmed,
                topics_json, competencies_json, created_at
              FROM {schema}.education_items
              WHERE user_id = %s AND archived_at IS NULL"""
    params = [user["id"]]

    if kind_filter == "formal":
        sql += " AND kind IN ('degree','certificate','course','program')"
    elif kind_filter == "material":
        sql += " AND kind IN ('book','lecture','presentation','methodology','notes','article','material')"
    elif kind_filter and kind_filter != "all":
        sql += " AND kind = %s"
        params.append(kind_filter)

    if status_filter and status_filter != "all":
        sql += " AND status = %s"
        params.append(status_filter)

    sql += " ORDER BY created_at DESC"

    cur.execute(sql, tuple(params))
    rows = cur.fetchall()
    items = []
    for r in rows:
        topics = []
        comps = []
        try:
            if r[12]:
                topics = json.loads(r[12])
        except Exception:
            pass
        try:
            if r[13]:
                comps = json.loads(r[13])
        except Exception:
            pass
        items.append({
            "id": r[0], "kind": r[1], "title": r[2],
            "issuer_name": r[3], "institution_name": r[4],
            "field_of_study": r[5], "level": r[6],
            "issued_at": str(r[7]) if r[7] else None,
            "status": r[8], "study_status": r[9],
            "source_type": r[10], "is_confirmed": r[11],
            "topics": topics, "competencies": comps,
            "created_at": str(r[14]),
        })
    return ok_response({"items": items, "total": len(items)}, request_id, origin)


def handle_get(conn, user, body, request_id, origin):
    schema = get_schema()
    item_id = body.get("id")
    if not item_id:
        return err_response("validation_error", "Нужен id", 400, request_id, origin)
    cur = conn.cursor()
    cur.execute(
        f"""SELECT id, user_id, kind, title, issuer_name, institution_name, description,
            field_of_study, level, start_date, end_date, issued_at, hours, grade, language,
            status, study_status, confidence, source_type, is_confirmed, confirmed_at,
            extracted_json, topics_json, competencies_json, created_at, updated_at
            FROM {schema}.education_items WHERE id = %s""",
        (int(item_id),),
    )
    r = cur.fetchone()
    if not r:
        return err_response("not_found", "Не найдено", 404, request_id, origin)
    if r[1] != user["id"]:
        return err_response("access_denied", "Нет доступа", 403, request_id, origin)

    # Файлы
    cur.execute(
        f"""SELECT id, original_name, mime_type, size_bytes, parse_status, created_at
            FROM {schema}.education_item_files WHERE education_item_id = %s""",
        (int(item_id),),
    )
    files = [
        {"id": fr[0], "name": fr[1], "mime": fr[2], "size": fr[3], "parse_status": fr[4], "created_at": str(fr[5])}
        for fr in cur.fetchall()
    ]

    def _parse_json(s):
        if not s:
            return None
        try:
            return json.loads(s)
        except Exception:
            return None

    return ok_response({
        "id": r[0], "kind": r[2], "title": r[3],
        "issuer_name": r[4], "institution_name": r[5], "description": r[6],
        "field_of_study": r[7], "level": r[8],
        "start_date": str(r[9]) if r[9] else None,
        "end_date": str(r[10]) if r[10] else None,
        "issued_at": str(r[11]) if r[11] else None,
        "hours": r[12], "grade": r[13], "language": r[14],
        "status": r[15], "study_status": r[16],
        "confidence": r[17], "source_type": r[18],
        "is_confirmed": r[19], "confirmed_at": str(r[20]) if r[20] else None,
        "extracted_data": _parse_json(r[21]),
        "topics": _parse_json(r[22]) or [],
        "competencies": _parse_json(r[23]) or [],
        "created_at": str(r[24]), "updated_at": str(r[25]),
        "files": files,
    }, request_id, origin)


def handle_create(conn, user, body, request_id, origin):
    schema = get_schema()
    kind = body.get("kind")
    title = (body.get("title") or "").strip()
    if not kind or not title:
        return err_response("validation_error", "Нужны kind и title", 400, request_id, origin)
    cur = conn.cursor()
    cur.execute(
        f"""INSERT INTO {schema}.education_items
            (user_id, kind, title, issuer_name, institution_name, description,
             field_of_study, level, start_date, end_date, issued_at, hours, grade,
             study_status, source_type, status, is_confirmed, confirmed_at)
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, 'manual', 'confirmed', true, NOW())
            RETURNING id""",
        (
            user["id"], kind, title,
            body.get("issuer_name"), body.get("institution_name"), body.get("description"),
            body.get("field_of_study"), body.get("level"),
            body.get("start_date") or None, body.get("end_date") or None, body.get("issued_at") or None,
            body.get("hours"), body.get("grade"),
            body.get("study_status"),
        ),
    )
    new_id = cur.fetchone()[0]
    conn.commit()
    notify_indexer("upsert", new_id)
    return ok_response({"id": new_id, "title": title, "kind": kind}, request_id, origin)


def handle_update(conn, user, body, request_id, origin):
    schema = get_schema()
    item_id = body.get("id")
    if not item_id:
        return err_response("validation_error", "Нужен id", 400, request_id, origin)
    cur = conn.cursor()
    cur.execute(f"SELECT user_id FROM {schema}.education_items WHERE id = %s", (int(item_id),))
    row = cur.fetchone()
    if not row:
        return err_response("not_found", "Не найдено", 404, request_id, origin)
    if row[0] != user["id"]:
        return err_response("access_denied", "Нет доступа", 403, request_id, origin)

    fields = ("title", "issuer_name", "institution_name", "description",
              "field_of_study", "level", "start_date", "end_date", "issued_at",
              "hours", "grade", "study_status", "status")
    sets = []
    params = []
    for f in fields:
        if f in body:
            sets.append(f"{f} = %s")
            val = body[f]
            if f in ("start_date", "end_date", "issued_at") and not val:
                val = None
            params.append(val)
    if sets:
        sets.append("updated_at = NOW()")
        params.append(int(item_id))
        cur.execute(f"UPDATE {schema}.education_items SET {', '.join(sets)} WHERE id = %s", tuple(params))
        conn.commit()
        notify_indexer("upsert", int(item_id))
    return ok_response({"ok": True}, request_id, origin)


def handle_archive(conn, user, body, request_id, origin):
    schema = get_schema()
    item_id = body.get("id")
    if not item_id:
        return err_response("validation_error", "Нужен id", 400, request_id, origin)
    cur = conn.cursor()
    cur.execute(f"SELECT user_id, title FROM {schema}.education_items WHERE id = %s", (int(item_id),))
    row = cur.fetchone()
    if not row:
        return err_response("not_found", "Не найдено", 404, request_id, origin)
    if row[0] != user["id"]:
        return err_response("access_denied", "Нет доступа", 403, request_id, origin)
    cur.execute(f"UPDATE {schema}.education_items SET archived_at = NOW() WHERE id = %s", (int(item_id),))
    conn.commit()
    notify_indexer("delete", int(item_id))
    return ok_response({"ok": True, "archived": True}, request_id, origin)


def handle_upload_file(conn, user, body, request_id, origin):
    """Прикрепляет файл к записи + автозапускает AI-анализ."""
    schema = get_schema()
    item_id = body.get("id")
    file_data_b64 = body.get("file_data")
    filename = body.get("filename", "file")
    mime = body.get("mime", "application/octet-stream")
    if not item_id or not file_data_b64:
        return err_response("validation_error", "Нужны id и file_data", 400, request_id, origin)

    cur = conn.cursor()
    cur.execute(f"SELECT user_id, kind FROM {schema}.education_items WHERE id = %s", (int(item_id),))
    row = cur.fetchone()
    if not row:
        return err_response("not_found", "Запись не найдена", 404, request_id, origin)
    if row[0] != user["id"]:
        return err_response("access_denied", "Нет доступа", 403, request_id, origin)
    kind = row[1]

    file_bytes = base64.b64decode(file_data_b64)
    s3_key = f"education/{user['id']}/{item_id}_{filename}"

    # Загружаем в S3
    try:
        s3 = get_s3()
        s3.put_object(Bucket="files", Key=s3_key, Body=file_bytes, ContentType=mime)
    except Exception as e:
        return err_response("storage_error", f"Не удалось сохранить файл: {e}", 500, request_id, origin)

    # Извлекаем текст. Чёткая обработка случаев:
    #  - текст есть и нормальный → done
    #  - пусто или очень мало (<50 символов) → empty_or_too_short
    #  - начинается с [Ошибка / [OCR: текст не найден → failed
    parsed_text = extract_text_from_file(file_bytes, mime)
    if not parsed_text:
        parse_status = "empty"
    elif parsed_text.startswith("[Ошибка") or parsed_text.startswith("[OCR"):
        parse_status = "failed"
    elif len(parsed_text.strip()) < 50:
        parse_status = "too_short"
    else:
        parse_status = "done"

    cur.execute(
        f"""INSERT INTO {schema}.education_item_files
            (education_item_id, s3_key, original_name, mime_type, size_bytes, parsed_text, parse_status)
            VALUES (%s, %s, %s, %s, %s, %s, %s) RETURNING id""",
        (int(item_id), s3_key, filename, mime, len(file_bytes), parsed_text[:50000], parse_status),
    )
    file_id = cur.fetchone()[0]

    # Автозапуск AI-анализа если текст извлёкся
    extracted = None
    if parse_status == "done":
        cur.execute(
            f"UPDATE {schema}.education_items SET status = 'processing', source_type = 'uploaded_file' WHERE id = %s",
            (int(item_id),),
        )
        conn.commit()

        if kind in KIND_GROUPS["formal"]:
            extracted = ai_extract_formal(parsed_text)
        else:
            extracted = ai_extract_material(parsed_text)

        # Сохраняем извлечённое
        topics = extracted.get("topics") or extracted.get("main_topics") or []
        competencies = extracted.get("suggested_competencies") or []

        cur.execute(
            f"""UPDATE {schema}.education_items SET
                extracted_json = %s,
                topics_json = %s,
                competencies_json = %s,
                status = 'needs_review',
                source_type = 'ai_extracted'
                WHERE id = %s""",
            (
                json.dumps(extracted, ensure_ascii=False),
                json.dumps(topics, ensure_ascii=False),
                json.dumps(competencies, ensure_ascii=False),
                int(item_id),
            ),
        )

    conn.commit()

    # Понятные предупреждения для UI
    warning = None
    if parse_status == "failed":
        warning = "Не удалось извлечь текст из файла (возможно повреждён или зашифрован). Проверьте файл или введите данные вручную."
    elif parse_status == "empty":
        warning = "Файл пустой — нет текста для анализа. Введите данные вручную."
    elif parse_status == "too_short":
        warning = f"Из файла извлечено мало текста ({len(parsed_text.strip())} символов). AI-извлечение может быть неполным — проверьте результат."

    return ok_response({
        "file_id": file_id,
        "parse_status": parse_status,
        "extracted": extracted,
        "text_length": len(parsed_text or ""),
        "raw_text_preview": (parsed_text or "")[:1000] if parse_status in ("failed", "empty", "too_short") else None,
        "warning": warning,
    }, request_id, origin)


def handle_confirm(conn, user, body, request_id, origin):
    """Пользователь подтверждает AI-извлечение, опционально с правками."""
    schema = get_schema()
    item_id = body.get("id")
    if not item_id:
        return err_response("validation_error", "Нужен id", 400, request_id, origin)
    cur = conn.cursor()
    cur.execute(f"SELECT user_id FROM {schema}.education_items WHERE id = %s", (int(item_id),))
    row = cur.fetchone()
    if not row:
        return err_response("not_found", "Не найдено", 404, request_id, origin)
    if row[0] != user["id"]:
        return err_response("access_denied", "Нет доступа", 403, request_id, origin)

    # Применяем правки (если пользователь скорректировал)
    overrides = body.get("overrides") or {}
    sets = ["status = 'confirmed'", "is_confirmed = true", "confirmed_at = NOW()", "updated_at = NOW()"]
    params = []
    for f in ("title", "issuer_name", "institution_name", "field_of_study", "level",
              "issued_at", "hours", "grade", "study_status"):
        if f in overrides:
            sets.append(f"{f} = %s")
            params.append(overrides[f] or None)
    if "topics" in overrides:
        sets.append("topics_json = %s")
        params.append(json.dumps(overrides["topics"], ensure_ascii=False))
    if "competencies" in overrides:
        sets.append("competencies_json = %s")
        params.append(json.dumps(overrides["competencies"], ensure_ascii=False))

    params.append(int(item_id))
    cur.execute(f"UPDATE {schema}.education_items SET {', '.join(sets)} WHERE id = %s", tuple(params))
    conn.commit()
    notify_indexer("upsert", int(item_id))
    return ok_response({"ok": True, "confirmed": True}, request_id, origin)


def handle_analyze(conn, user, body, request_id, origin):
    """Повторный AI-анализ — перебирает первый файл записи."""
    schema = get_schema()
    item_id = body.get("id")
    if not item_id:
        return err_response("validation_error", "Нужен id", 400, request_id, origin)
    cur = conn.cursor()
    cur.execute(f"SELECT user_id, kind FROM {schema}.education_items WHERE id = %s", (int(item_id),))
    row = cur.fetchone()
    if not row:
        return err_response("not_found", "Не найдено", 404, request_id, origin)
    if row[0] != user["id"]:
        return err_response("access_denied", "Нет доступа", 403, request_id, origin)
    kind = row[1]

    cur.execute(
        f"SELECT parsed_text FROM {schema}.education_item_files WHERE education_item_id = %s ORDER BY id LIMIT 1",
        (int(item_id),),
    )
    fr = cur.fetchone()
    if not fr or not fr[0]:
        return err_response("no_file", "Нет файла или текст не извлечён", 400, request_id, origin)

    extracted = ai_extract_formal(fr[0]) if kind in KIND_GROUPS["formal"] else ai_extract_material(fr[0])
    topics = extracted.get("topics") or extracted.get("main_topics") or []
    competencies = extracted.get("suggested_competencies") or []

    cur.execute(
        f"""UPDATE {schema}.education_items SET
            extracted_json = %s, topics_json = %s, competencies_json = %s,
            status = 'needs_review', updated_at = NOW()
            WHERE id = %s""",
        (
            json.dumps(extracted, ensure_ascii=False),
            json.dumps(topics, ensure_ascii=False),
            json.dumps(competencies, ensure_ascii=False),
            int(item_id),
        ),
    )
    conn.commit()
    return ok_response({"extracted": extracted}, request_id, origin)


def handle_profile_summary(conn, user, request_id, origin):
    """Сводка для Dashboard: счётчики + области знаний."""
    schema = get_schema()
    cur = conn.cursor()

    # Счётчики по типам
    cur.execute(
        f"""SELECT kind, COUNT(*)
            FROM {schema}.education_items
            WHERE user_id = %s AND archived_at IS NULL
            GROUP BY kind""",
        (user["id"],),
    )
    counts = {r[0]: r[1] for r in cur.fetchall()}

    formal_count = sum(counts.get(k, 0) for k in ("degree", "certificate", "course", "program"))
    material_count = sum(counts.get(k, 0) for k in ("book", "lecture", "presentation", "methodology", "notes", "article", "material"))

    # Сколько подтверждённых
    cur.execute(
        f"""SELECT COUNT(*) FROM {schema}.education_items
            WHERE user_id = %s AND archived_at IS NULL AND is_confirmed = true""",
        (user["id"],),
    )
    confirmed = cur.fetchone()[0]

    # Сколько ждут проверки
    cur.execute(
        f"""SELECT COUNT(*) FROM {schema}.education_items
            WHERE user_id = %s AND archived_at IS NULL AND status = 'needs_review'""",
        (user["id"],),
    )
    needs_review = cur.fetchone()[0]

    # Области знаний — собираем все topics
    cur.execute(
        f"""SELECT topics_json FROM {schema}.education_items
            WHERE user_id = %s AND archived_at IS NULL AND topics_json IS NOT NULL""",
        (user["id"],),
    )
    domain_counter = {}
    for (tj,) in cur.fetchall():
        try:
            topics = json.loads(tj)
            for t in topics:
                t_clean = (t or "").strip().lower()
                if t_clean:
                    domain_counter[t_clean] = domain_counter.get(t_clean, 0) + 1
        except Exception:
            continue
    top_topics = sorted(domain_counter.items(), key=lambda x: -x[1])[:10]

    return ok_response({
        "degree": counts.get("degree", 0),
        "certificate": counts.get("certificate", 0),
        "course": counts.get("course", 0),
        "program": counts.get("program", 0),
        "material_total": material_count,
        "formal_total": formal_count,
        "confirmed_total": confirmed,
        "needs_review_total": needs_review,
        "top_topics": [{"name": t, "count": c} for t, c in top_topics],
    }, request_id, origin)


# ============================================================
# Presigned upload — файл заливается напрямую из браузера в S3
# ============================================================


def handle_get_upload_url(conn, user, body, request_id, origin):
    """Возвращает presigned PUT URL для загрузки файла напрямую в S3 из браузера.
    Браузер НЕ шлёт base64 через нас — сразу PUT на S3 (обходит лимит 1МБ)."""
    schema = get_schema()
    item_id = body.get("id")
    filename = body.get("filename", "file")
    mime = body.get("mime", "application/octet-stream")
    if not item_id or not filename:
        return err_response("validation_error", "Нужны id и filename", 400, request_id, origin)

    cur = conn.cursor()
    cur.execute(f"SELECT user_id FROM {schema}.education_items WHERE id = %s", (int(item_id),))
    row = cur.fetchone()
    if not row:
        return err_response("not_found", "Запись не найдена", 404, request_id, origin)
    if row[0] != user["id"]:
        return err_response("access_denied", "Нет доступа", 403, request_id, origin)

    import re
    safe_name = re.sub(r"[^\w.\-]", "_", filename)[:120]
    s3_key = f"education/{user['id']}/{item_id}_{safe_name}"

    s3 = get_s3()
    presigned_url = s3.generate_presigned_url(
        "put_object",
        Params={"Bucket": "files", "Key": s3_key, "ContentType": mime},
        ExpiresIn=600,
    )
    return ok_response({"upload_url": presigned_url, "s3_key": s3_key}, request_id, origin)


def handle_file_ready(conn, user, body, request_id, origin):
    """Вызывается после того как фронт залил файл в S3 через presigned URL.
    Создаёт запись в education_item_files, скачивает файл из S3, парсит, запускает AI."""
    schema = get_schema()
    item_id = body.get("id")
    filename = body.get("filename", "file")
    mime = body.get("mime", "application/octet-stream")
    s3_key = body.get("s3_key")
    file_size = body.get("file_size", 0)
    if not item_id or not s3_key:
        return err_response("validation_error", "Нужны id и s3_key", 400, request_id, origin)

    cur = conn.cursor()
    cur.execute(f"SELECT user_id, kind FROM {schema}.education_items WHERE id = %s", (int(item_id),))
    row = cur.fetchone()
    if not row:
        return err_response("not_found", "Запись не найдена", 404, request_id, origin)
    if row[0] != user["id"]:
        return err_response("access_denied", "Нет доступа", 403, request_id, origin)
    kind = row[1]

    # Скачиваем файл из S3 для парсинга
    s3 = get_s3()
    file_bytes = b""
    try:
        obj = s3.get_object(Bucket="files", Key=s3_key)
        file_bytes = obj["Body"].read()
    except Exception as e:
        log.warning("Не удалось скачать файл из S3: %s", e)

    # Парсим текст
    parsed_text = extract_text_from_file(file_bytes, mime) if file_bytes else ""
    if not parsed_text:
        parse_status = "empty"
    elif parsed_text.startswith("[Ошибка") or parsed_text.startswith("[OCR"):
        parse_status = "failed"
    elif len(parsed_text.strip()) < 50:
        parse_status = "too_short"
    else:
        parse_status = "done"

    cur.execute(
        f"""INSERT INTO {schema}.education_item_files
            (education_item_id, s3_key, original_name, mime_type, size_bytes, parsed_text, parse_status)
            VALUES (%s, %s, %s, %s, %s, %s, %s) RETURNING id""",
        (int(item_id), s3_key, filename, mime, file_size or len(file_bytes), parsed_text[:50000], parse_status),
    )
    file_id = cur.fetchone()[0]

    extracted = None
    if parse_status == "done":
        cur.execute(
            f"UPDATE {schema}.education_items SET status = 'processing', source_type = 'uploaded_file' WHERE id = %s",
            (int(item_id),),
        )
        conn.commit()
        extracted = ai_extract_formal(parsed_text) if kind in KIND_GROUPS["formal"] else ai_extract_material(parsed_text)
        topics = extracted.get("topics") or extracted.get("main_topics") or []
        competencies = extracted.get("suggested_competencies") or []
        cur.execute(
            f"""UPDATE {schema}.education_items SET
                extracted_json = %s, topics_json = %s, competencies_json = %s,
                status = 'needs_review', source_type = 'ai_extracted', updated_at = NOW()
                WHERE id = %s""",
            (json.dumps(extracted, ensure_ascii=False), json.dumps(topics, ensure_ascii=False),
             json.dumps(competencies, ensure_ascii=False), int(item_id)),
        )

    warning = None
    if parse_status == "failed":
        warning = "Не удалось извлечь текст (файл повреждён или зашифрован). Введите данные вручную."
    elif parse_status == "empty":
        warning = "Файл пустой. Введите данные вручную."
    elif parse_status == "too_short":
        warning = f"Извлечено мало текста ({len(parsed_text.strip())} симв.). Проверьте результат AI."

    conn.commit()
    return ok_response({
        "file_id": file_id, "parse_status": parse_status,
        "extracted": extracted, "warning": warning,
    }, request_id, origin)


def handle_get_file_url(conn, user, body, request_id, origin):
    """Возвращает presigned GET URL для скачивания/просмотра файла."""
    schema = get_schema()
    file_id = body.get("file_id")
    if not file_id:
        return err_response("validation_error", "Нужен file_id", 400, request_id, origin)

    cur = conn.cursor()
    cur.execute(
        f"""SELECT eif.s3_key, eif.original_name, eif.mime_type, ei.user_id
            FROM {schema}.education_item_files eif
            JOIN {schema}.education_items ei ON ei.id = eif.education_item_id
            WHERE eif.id = %s""",
        (int(file_id),),
    )
    row = cur.fetchone()
    if not row:
        return err_response("not_found", "Файл не найден", 404, request_id, origin)
    if row[3] != user["id"]:
        return err_response("access_denied", "Нет доступа", 403, request_id, origin)

    s3_key, original_name, mime_type = row[0], row[1], row[2]
    s3 = get_s3()
    try:
        url = s3.generate_presigned_url(
            "get_object",
            Params={"Bucket": "files", "Key": s3_key},
            ExpiresIn=300,
        )
    except Exception as e:
        return err_response("storage_error", f"Не удалось получить ссылку: {e}", 500, request_id, origin)

    return ok_response({"url": url, "filename": original_name, "mime": mime_type}, request_id, origin)


# ============================================================
# Router
# ============================================================


def handler(event: dict, context) -> dict:
    request_id = getattr(context, "request_id", None) or str(uuid.uuid4())
    origin = event.get("headers", {}).get("Origin") or event.get("headers", {}).get("origin")

    if event.get("httpMethod") == "OPTIONS":
        return {"statusCode": 200, "headers": cors_headers(origin), "body": ""}

    if event.get("httpMethod") != "POST":
        return err_response("method_not_allowed", "Используйте POST", 405, request_id, origin)

    body = {}
    if event.get("body"):
        try:
            body = json.loads(event["body"])
        except Exception:
            return err_response("invalid_json", "Невалидный JSON", 400, request_id, origin)

    action = body.get("action", "")
    if action not in ALLOWED_ACTIONS:
        return err_response("unknown_action", f"Допустимые: {sorted(ALLOWED_ACTIONS)}", 400, request_id, origin)

    session_id = event.get("headers", {}).get("X-Session-Id", "")
    conn = get_db()
    try:
        user = get_current_user(conn, session_id)
        if not user:
            return err_response("auth_required", "Требуется авторизация", 401, request_id, origin)

        log.info("request_id=%s action=%s user=%s", request_id, action, user["id"])

        if action == "education.list":
            return handle_list(conn, user, body, request_id, origin)
        if action == "education.get":
            return handle_get(conn, user, body, request_id, origin)
        if action == "education.create":
            return handle_create(conn, user, body, request_id, origin)
        if action == "education.update":
            return handle_update(conn, user, body, request_id, origin)
        if action == "education.archive":
            return handle_archive(conn, user, body, request_id, origin)
        if action == "education.upload_file":
            return handle_upload_file(conn, user, body, request_id, origin)
        if action == "education.get_upload_url":
            return handle_get_upload_url(conn, user, body, request_id, origin)
        if action == "education.file_ready":
            return handle_file_ready(conn, user, body, request_id, origin)
        if action == "education.get_file_url":
            return handle_get_file_url(conn, user, body, request_id, origin)
        if action == "education.analyze":
            return handle_analyze(conn, user, body, request_id, origin)
        if action == "education.confirm":
            return handle_confirm(conn, user, body, request_id, origin)
        if action == "education.profile_summary":
            return handle_profile_summary(conn, user, request_id, origin)

        return err_response("not_implemented", "Не реализовано", 501, request_id, origin)

    except Exception as e:
        log.exception("Unhandled error request_id=%s", request_id)
        return err_response("internal_error", f"Ошибка: {str(e)[:200]}", 500, request_id, origin)
    finally:
        conn.close()