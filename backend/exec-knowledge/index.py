import json
import os
import io
import uuid
import hashlib
import base64
import psycopg2
import boto3

DB = os.environ["DATABASE_URL"]
_s = os.environ.get("MAIN_DB_SCHEMA", "").strip()
SCHEMA = _s if _s else "t_p61016064_digital_innovation_i"

MAX_TEXT_LEN = 400000
CHUNK_SIZE = 1500
CHUNK_OVERLAP = 200

DOC_TYPES = {
    "rule": "Регламент",
    "matrix": "Матрица ответственности",
    "policy": "Политика",
    "method": "Методика",
    "template": "Шаблон",
    "note": "Вводная",
    "other": "Другое",
}


def cors(body: dict, code: int = 200) -> dict:
    return {
        "statusCode": code,
        "headers": {
            "Access-Control-Allow-Origin": "*",
            "Access-Control-Allow-Methods": "GET, POST, OPTIONS",
            "Access-Control-Allow-Headers": "Content-Type, X-Admin-Token, X-Session-Id",
            "Content-Type": "application/json",
        },
        "body": json.dumps(body, ensure_ascii=False, default=str),
    }


def get_admin(conn, token: str):
    if not token:
        return None
    token_hash = hashlib.sha256(token.encode()).hexdigest()
    with conn.cursor() as cur:
        cur.execute(
            f"SELECT actor_email FROM {SCHEMA}.admin_sessions "
            f"WHERE session_token_hash = %s AND expires_at > NOW() AND revoked_at IS NULL LIMIT 1",
            (token_hash,),
        )
        row = cur.fetchone()
    return row[0] if row else None


def get_cabinet_user(conn, session_id: str):
    if not session_id:
        return None
    with conn.cursor() as cur:
        cur.execute(
            f"SELECT u.email, a.access_role, a.can_confirm "
            f"FROM {SCHEMA}.sessions s "
            f"JOIN {SCHEMA}.users u ON u.id = s.user_id "
            f"JOIN {SCHEMA}.exec_cabinet_access a ON LOWER(a.email) = LOWER(u.email) "
            f"LEFT JOIN {SCHEMA}.admin_user_flags fl ON fl.user_id = u.id "
            f"WHERE s.id = %s AND s.expires_at > NOW() "
            f"AND a.is_active = true AND COALESCE(fl.is_blocked, false) = false LIMIT 1",
            (session_id,),
        )
        row = cur.fetchone()
    if not row:
        return None
    return {"email": row[0], "role": row[1], "can_confirm": row[2]}


def authenticate(conn, headers: dict):
    token = headers.get("x-admin-token") or headers.get("X-Admin-Token", "")
    email = get_admin(conn, token)
    if email:
        return {"email": email, "role": "head", "can_confirm": True}
    sid = headers.get("x-session-id") or headers.get("X-Session-Id", "")
    return get_cabinet_user(conn, sid)


def rows(cur):
    cols = [d[0] for d in cur.description]
    return [dict(zip(cols, r)) for r in cur.fetchall()]


def nz(v):
    if v is None:
        return None
    if isinstance(v, str) and not v.strip():
        return None
    return v


def as_int(v):
    v = nz(v)
    try:
        return int(v) if v is not None else None
    except Exception:
        return None


def get_s3():
    return boto3.client(
        "s3",
        endpoint_url="https://bucket.poehali.dev",
        aws_access_key_id=os.environ["AWS_ACCESS_KEY_ID"],
        aws_secret_access_key=os.environ["AWS_SECRET_ACCESS_KEY"],
    )


def safe_name(filename: str) -> str:
    import re
    table = str.maketrans({
        "а": "a", "б": "b", "в": "v", "г": "g", "д": "d", "е": "e", "ё": "e",
        "ж": "zh", "з": "z", "и": "i", "й": "y", "к": "k", "л": "l", "м": "m",
        "н": "n", "о": "o", "п": "p", "р": "r", "с": "s", "т": "t", "у": "u",
        "ф": "f", "х": "h", "ц": "ts", "ч": "ch", "ш": "sh", "щ": "sch",
        "ъ": "", "ы": "y", "ь": "", "э": "e", "ю": "yu", "я": "ya",
    })
    name = (filename or "file").lower().translate(table)
    name = re.sub(r"[^a-z0-9._-]+", "_", name)
    name = re.sub(r"_+", "_", name).strip("_")
    return name[:150] or "file"


def chunk_text(text: str) -> list:
    """Делит текст на перекрывающиеся фрагменты."""
    chunks = []
    if not text:
        return chunks
    start, idx = 0, 0
    while start < len(text):
        end = min(start + CHUNK_SIZE, len(text))
        if end < len(text):
            for sep in [". ", "!\n", "?\n", "\n\n", ".\n", "\n"]:
                cut = text.rfind(sep, start, end)
                if cut > start + CHUNK_SIZE // 2:
                    end = cut + len(sep)
                    break
        piece = text[start:end].strip()
        if piece:
            chunks.append({"index": idx, "content": piece, "page": None})
            idx += 1
        start = end - CHUNK_OVERLAP if end - CHUNK_OVERLAP > start else end
    return chunks


def extract_pdf(data: bytes):
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(data))
        pages = []
        for i, page in enumerate(reader.pages):
            pages.append((i + 1, page.extract_text() or ""))
        text = "\n".join(p[1] for p in pages)[:MAX_TEXT_LEN]
        chunks, idx = [], 0
        for page_num, page_text in pages:
            if not page_text.strip():
                continue
            for c in chunk_text(page_text):
                chunks.append({"index": idx, "content": c["content"], "page": page_num})
                idx += 1
        return text, chunks, len(reader.pages)
    except Exception as e:
        return f"[Ошибка чтения PDF: {e}]", [], None


def extract_docx(data: bytes):
    try:
        import docx
        doc = docx.Document(io.BytesIO(data))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        # таблицы важны: матрицы ответственности часто именно в них
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        text = "\n".join(parts)[:MAX_TEXT_LEN]
        return text, chunk_text(text), None
    except Exception as e:
        return f"[Ошибка чтения DOCX: {e}]", [], None


def extract_pptx(data: bytes):
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides):
            lines = [f"Слайд {i + 1}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lines.append(t)
            parts.append("\n".join(lines))
        text = "\n\n".join(parts)[:MAX_TEXT_LEN]
        return text, chunk_text(text), len(prs.slides)
    except Exception as e:
        return f"[Ошибка чтения PPTX: {e}]", [], None


def save_chunks(cur, knowledge_id: int, chunks: list):
    cur.execute(
        f"DELETE FROM {SCHEMA}.exec_knowledge_chunk WHERE knowledge_id = %s",
        (knowledge_id,),
    )
    for c in chunks:
        cur.execute(
            f"INSERT INTO {SCHEMA}.exec_knowledge_chunk "
            f"(knowledge_id, chunk_index, page_number, content, content_length) "
            f"VALUES (%s,%s,%s,%s,%s)",
            (knowledge_id, c["index"], c.get("page"), c["content"], len(c["content"])),
        )


def list_items(cur):
    cur.execute(f"""
        SELECT k.id, k.title, k.doc_type, k.summary, k.filename, k.file_type,
               k.file_size, k.page_count, k.extracted_length, k.use_in_ai,
               k.priority, k.status, k.created_by, k.created_at, k.updated_at,
               (k.body IS NOT NULL AND length(k.body) > 0) AS has_text,
               (SELECT COUNT(*) FROM {SCHEMA}.exec_knowledge_chunk c
                 WHERE c.knowledge_id = k.id) AS chunks
        FROM {SCHEMA}.exec_knowledge k
        WHERE k.status <> 'archived'
        ORDER BY k.use_in_ai DESC, k.priority DESC, k.created_at DESC
    """)
    return rows(cur)


def get_item(cur, item_id: int):
    cur.execute(f"SELECT * FROM {SCHEMA}.exec_knowledge WHERE id = %s", (item_id,))
    got = rows(cur)
    return got[0] if got else None


def save_note(cur, body: dict, actor: str):
    """Ручная запись: правило, вводная, выдержка из регламента."""
    item_id = as_int(body.get("id"))
    title = (body.get("title") or "").strip()
    if not title:
        return None, "Укажите название"
    text = (body.get("body") or "").strip()

    vals = {
        "title": title[:500],
        "doc_type": body.get("doc_type") if body.get("doc_type") in DOC_TYPES else "rule",
        "summary": nz(body.get("summary")),
        "body": nz(text),
        "use_in_ai": bool(body.get("use_in_ai", True)),
        "priority": as_int(body.get("priority")) or 50,
        "extracted_length": len(text) if text else 0,
    }

    if item_id:
        sets = ", ".join(f"{k} = %s" for k in vals)
        cur.execute(
            f"UPDATE {SCHEMA}.exec_knowledge SET {sets}, updated_at = now() "
            f"WHERE id = %s RETURNING id",
            list(vals.values()) + [item_id],
        )
    else:
        vals["created_by"] = actor
        cols = ", ".join(vals)
        ph = ", ".join(["%s"] * len(vals))
        cur.execute(
            f"INSERT INTO {SCHEMA}.exec_knowledge ({cols}) VALUES ({ph}) RETURNING id",
            list(vals.values()),
        )
    new_id = cur.fetchone()[0]

    if text:
        save_chunks(cur, new_id, chunk_text(text))
    return new_id, None


def upload_file(cur, body: dict, actor: str):
    """Загрузка файла: PDF, DOCX, PPTX, TXT."""
    filename = (body.get("filename") or "").strip()
    file_b64 = body.get("file_data") or ""
    if not filename or not file_b64:
        return None, "Не хватает данных файла"

    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in ("pdf", "docx", "pptx", "txt", "md"):
        return None, "Поддерживаются PDF, DOCX, PPTX, TXT"

    try:
        data = base64.b64decode(file_b64)
    except Exception:
        return None, "Файл повреждён при передаче"

    if len(data) > 20 * 1024 * 1024:
        return None, "Файл больше 20 МБ — загрузите документ поменьше"

    if ext == "pdf":
        text, chunks, pages = extract_pdf(data)
    elif ext == "docx":
        text, chunks, pages = extract_docx(data)
    elif ext == "pptx":
        text, chunks, pages = extract_pptx(data)
    else:
        text = data.decode("utf-8", errors="replace")[:MAX_TEXT_LEN]
        chunks, pages = chunk_text(text), None

    s3_key = None
    try:
        s3_key = f"exec-knowledge/{uuid.uuid4().hex}_{safe_name(filename)}"
        get_s3().put_object(Bucket="files", Key=s3_key, Body=data)
    except Exception:
        s3_key = None

    title = (body.get("title") or "").strip() or filename.rsplit(".", 1)[0]
    cur.execute(
        f"INSERT INTO {SCHEMA}.exec_knowledge "
        f"(title, doc_type, summary, body, filename, file_type, file_size, s3_key, "
        f" page_count, extracted_length, use_in_ai, priority, created_by) "
        f"VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s) RETURNING id",
        (title[:500],
         body.get("doc_type") if body.get("doc_type") in DOC_TYPES else "rule",
         nz(body.get("summary")), text, filename[:255], ext, len(data), s3_key,
         pages, len(text), bool(body.get("use_in_ai", True)),
         as_int(body.get("priority")) or 50, actor),
    )
    new_id = cur.fetchone()[0]
    save_chunks(cur, new_id, chunks)
    return {"id": new_id, "chunks": len(chunks), "length": len(text)}, None


def handler(event: dict, context) -> dict:
    """База знаний кабинета руководителя: регламенты и правила для AI-контекста."""
    if event.get("httpMethod") == "OPTIONS":
        return cors({})

    headers = event.get("headers") or {}
    conn = psycopg2.connect(DB)
    try:
        user = authenticate(conn, headers)
        if not user:
            return cors({"ok": False, "error": {"message": "Не авторизован"}}, 401)

        qs = event.get("queryStringParameters") or {}
        action = qs.get("action", "list")
        body = json.loads(event["body"]) if event.get("body") else {}
        cur = conn.cursor()

        if action == "list":
            return cors({"ok": True, "data": {
                "items": list_items(cur),
                "doc_types": DOC_TYPES,
            }})

        if action == "get":
            item_id = as_int(qs.get("id")) or as_int(body.get("id"))
            if not item_id:
                return cors({"ok": False, "error": {"message": "Не указан документ"}}, 400)
            item = get_item(cur, item_id)
            if not item:
                return cors({"ok": False, "error": {"message": "Документ не найден"}}, 404)
            return cors({"ok": True, "data": item})

        if action == "save_note":
            new_id, err = save_note(cur, body, user["email"])
            if err:
                return cors({"ok": False, "error": {"message": err}}, 400)
            conn.commit()
            return cors({"ok": True, "data": {"id": new_id}})

        if action == "upload":
            data, err = upload_file(cur, body, user["email"])
            if err:
                return cors({"ok": False, "error": {"message": err}}, 400)
            conn.commit()
            return cors({"ok": True, "data": data})

        if action == "toggle_ai":
            item_id = as_int(body.get("id"))
            if not item_id:
                return cors({"ok": False, "error": {"message": "Не указан документ"}}, 400)
            cur.execute(
                f"UPDATE {SCHEMA}.exec_knowledge SET use_in_ai = %s, updated_at = now() "
                f"WHERE id = %s RETURNING use_in_ai",
                (bool(body.get("use_in_ai")), item_id),
            )
            row = cur.fetchone()
            conn.commit()
            return cors({"ok": True, "data": {"id": item_id, "use_in_ai": row[0] if row else None}})

        if action == "delete":
            item_id = as_int(body.get("id"))
            if not item_id:
                return cors({"ok": False, "error": {"message": "Не указан документ"}}, 400)
            cur.execute(
                f"UPDATE {SCHEMA}.exec_knowledge SET status = 'archived', "
                f"use_in_ai = false, updated_at = now() WHERE id = %s",
                (item_id,),
            )
            conn.commit()
            return cors({"ok": True, "data": {"id": item_id}})

        return cors({"ok": False, "error": {"message": f"Неизвестное действие: {action}"}}, 400)
    finally:
        conn.close()
