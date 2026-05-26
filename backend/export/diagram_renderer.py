"""
Diagram Renderer — рисует схемы через python-pptx shapes в заданной зоне слайда.

Все renderer-функции принимают (prompt, slide, style, zx, zy, zw, zh):
  zx, zy — верхний левый угол зоны (Inches от края слайда)
  zw, zh — ширина и высота зоны (Inches)

Координаты всех shapes вычисляются с учётом смещения зоны.
"""
import re
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor


# ------------------------------------------------------------------ #
#  Цвет-утилиты                                                       #
# ------------------------------------------------------------------ #

def _rgb(hex_str: str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _darken(hex_color: str, factor: float) -> str:
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return "#{:02X}{:02X}{:02X}".format(
        max(0, int(r * (1 - factor))),
        max(0, int(g * (1 - factor))),
        max(0, int(b * (1 - factor))),
    )


def _lighten(hex_color: str, factor: float) -> str:
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return "#{:02X}{:02X}{:02X}".format(
        min(255, int(r + (255 - r) * factor)),
        min(255, int(g + (255 - g) * factor)),
        min(255, int(b + (255 - b) * factor)),
    )


# ------------------------------------------------------------------ #
#  Парсинг шагов                                                      #
# ------------------------------------------------------------------ #

def _parse_steps(prompt: str, max_steps: int = 8) -> list:
    """Извлекает список шагов. Поддерживает →, ->, em-dash, запятые, нумерацию."""
    norm = prompt.replace("→", "->").replace("—", ",")

    if "->" in norm:
        parts = norm.split("->")
        cleaned = []
        for s in parts:
            s = s.strip().strip(".,;")
            if ":" in s:
                prefix, _, rest = s.partition(":")
                if len(prefix.strip()) < 40:
                    s = rest.strip()
            if s:
                cleaned.append(s[:55])
        if len(cleaned) >= 2:
            return cleaned[:max_steps]

    numbered = re.findall(r'(?:^|\b)(\d+)[.\)]\s*([^\d\n;]+?)(?=\s*\d+[.\)]|\Z)', norm)
    if len(numbered) >= 2:
        return [v.strip().strip(",")[:55] for _, v in numbered][:max_steps]

    parts = [p.strip()[:55] for p in norm.split(",") if p.strip()]
    if len(parts) >= 2:
        return parts[:max_steps]

    return [prompt[:55]]


# ------------------------------------------------------------------ #
#  Shape-примитивы (координаты — абсолютные Inches)                  #
# ------------------------------------------------------------------ #

def _add_rect(slide, x, y, w, h, fill_hex, text, font_size=10, text_hex="#FFFFFF", bold=True):
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.color.rgb = _rgb(fill_hex)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(text_hex)
    return shape


def _add_connector(slide, x1, y1, x2, y2, color_hex):
    c = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    c.line.color.rgb = _rgb(color_hex)
    c.line.width = Emu(20000)
    return c


def _add_label(slide, x, y, w, h, text, font_size, color_hex, bold=False, align="center"):
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return tb


# ------------------------------------------------------------------ #
#  Renderers                                                          #
# ------------------------------------------------------------------ #

def _render_process(prompt, slide, style, zx, zy, zw, zh):
    """Горизонтальная цепочка блоков со стрелками внутри зоны."""
    steps = _parse_steps(prompt)
    n = len(steps)
    if not n:
        return False

    acc = "#{:02X}{:02X}{:02X}".format(*style.get("accent", (0x60, 0x8B, 0xC4)))
    txt = "#{:02X}{:02X}{:02X}".format(*style.get("title", (0xFF, 0xFF, 0xFF)))

    pad_x = 0.12
    arrow_w = 0.28
    avail_w = zw - pad_x * 2 - arrow_w * (n - 1)
    bw = min(1.7, avail_w / max(n, 1))
    bh = min(1.3, zh * 0.42)
    by = zy + (zh - bh) / 2

    total_used = bw * n + arrow_w * (n - 1)
    start_x = zx + pad_x + (zw - pad_x * 2 - total_used) / 2
    fs = max(8, min(11, int(13 - n * 0.5)))

    for i, step in enumerate(steps):
        x = start_x + i * (bw + arrow_w)
        fill = acc if i % 2 == 0 else _darken(acc, 0.22)
        _add_rect(slide, x, by, bw, bh, fill, step,
                  font_size=fs, text_hex=txt, bold=True)
        if i < n - 1:
            _add_connector(slide,
                           x + bw + 0.02, by + bh / 2,
                           x + bw + arrow_w - 0.04, by + bh / 2,
                           acc)
    return True


def _render_timeline(prompt, slide, style, zx, zy, zw, zh):
    """Горизонтальная ось с точками и метками (чередование вверх/вниз) внутри зоны."""
    from pptx.enum.shapes import MSO_SHAPE

    steps = _parse_steps(prompt, max_steps=8)
    n = len(steps)
    if not n:
        return False

    acc = "#{:02X}{:02X}{:02X}".format(*style.get("accent", (0x60, 0x8B, 0xC4)))
    txt = "#{:02X}{:02X}{:02X}".format(*style.get("title", (0xFF, 0xFF, 0xFF)))

    pad_x = 0.4
    axis_y = zy + zh * 0.48
    axis_x1 = zx + pad_x
    axis_x2 = zx + zw - pad_x

    # Горизонтальная ось
    ax = slide.shapes.add_connector(
        1,
        Inches(axis_x1), Inches(axis_y),
        Inches(axis_x2), Inches(axis_y),
    )
    ax.line.color.rgb = _rgb(acc)
    ax.line.width = Emu(30000)

    sg = (axis_x2 - axis_x1) / max(n - 1, 1) if n > 1 else (axis_x2 - axis_x1) / 2
    dot_r = 0.09
    label_arm = min(1.0, zh * 0.32)
    label_h = min(0.7, zh * 0.22)
    label_w = min(1.5, max(0.8, sg * 0.9))
    fs = max(7, min(9, int(10 - n * 0.25)))

    for i, step in enumerate(steps):
        sx = axis_x1 + i * sg

        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(sx - dot_r), Inches(axis_y - dot_r),
            Inches(dot_r * 2), Inches(dot_r * 2),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(acc)
        dot.line.fill.background()

        if i % 2 == 0:
            ty = axis_y - label_arm - label_h
            line_y1 = axis_y - dot_r
            line_y2 = ty + label_h
        else:
            ty = axis_y + label_arm
            line_y1 = axis_y + dot_r
            line_y2 = ty

        _add_connector(slide, sx, line_y1, sx, line_y2, acc)

        # Не выходим за границы зоны
        lx = max(zx + 0.05, min(sx - label_w / 2, zx + zw - label_w - 0.05))
        # Не выходим за верхнюю/нижнюю границу
        ty = max(zy + 0.05, min(ty, zy + zh - label_h - 0.05))
        _add_label(slide, lx, ty, label_w, label_h, step,
                   fs, txt, bold=(i in (0, n - 1)))

    return True


def _render_comparison(prompt, slide, style, zx, zy, zw, zh):
    """Два столбца VS внутри зоны."""
    from pptx.enum.text import PP_ALIGN

    acc = "#{:02X}{:02X}{:02X}".format(*style.get("accent", (0x60, 0x8B, 0xC4)))
    txt = "#{:02X}{:02X}{:02X}".format(*style.get("title", (0xFF, 0xFF, 0xFF)))

    parts = re.split(r'\bи\b|vs\.?|versus|\bпротив\b', prompt, flags=re.IGNORECASE)
    left = parts[0].strip()[:40] if len(parts) >= 2 else "Вариант А"
    right = parts[1].strip()[:40] if len(parts) >= 2 else "Вариант Б"

    vs_w = 0.6
    pad = 0.15
    cw = (zw - vs_w - pad * 4) / 2
    ch = zh * 0.72
    cy = zy + (zh - ch) / 2

    _add_rect(slide, zx + pad, cy, cw, ch,
              acc, left, font_size=11, text_hex=txt, bold=True)
    _add_rect(slide, zx + pad + cw + vs_w + pad * 2, cy, cw, ch,
              _darken(acc, 0.25), right, font_size=11, text_hex=txt, bold=True)
    _add_label(slide,
               zx + pad + cw + pad * 0.5, cy + ch / 2 - 0.22,
               vs_w, 0.44, "VS", 14, txt, bold=True)
    return True


def _render_matrix(prompt, slide, style, zx, zy, zw, zh):
    """2×2 матрица внутри зоны."""
    acc = "#{:02X}{:02X}{:02X}".format(*style.get("accent", (0x60, 0x8B, 0xC4)))
    txt = "#{:02X}{:02X}{:02X}".format(*style.get("title", (0xFF, 0xFF, 0xFF)))

    colors = [_darken(acc, 0.3), acc, _lighten(acc, 0.3), _darken(acc, 0.12)]
    labels = ["Низкий / Низкий", "Высокий / Низкий", "Низкий / Высокий", "Высокий / Высокий"]

    gap = 0.08
    cell_w = (zw - gap * 3) / 2
    cell_h = (zh - gap * 3) / 2

    for idx in range(4):
        col = idx % 2
        row = idx // 2
        x = zx + gap + col * (cell_w + gap)
        y = zy + gap + row * (cell_h + gap)
        fs = max(7, int(9 - (cell_w < 1.5) * 1))
        _add_rect(slide, x, y, cell_w, cell_h,
                  colors[idx], labels[idx], font_size=fs, text_hex=txt, bold=False)
    return True


def _render_orgchart(prompt, slide, style, zx, zy, zw, zh):
    """Дерево: корень + дочерние узлы."""
    acc = "#{:02X}{:02X}{:02X}".format(*style.get("accent", (0x60, 0x8B, 0xC4)))
    txt = "#{:02X}{:02X}{:02X}".format(*style.get("title", (0xFF, 0xFF, 0xFF)))

    steps = _parse_steps(prompt)
    if not steps:
        return False

    root = steps[0]
    children = steps[1:] if len(steps) > 1 else ["Команда"]

    rw = min(zw * 0.55, 2.4)
    rh = min(zh * 0.2, 0.8)
    rx = zx + (zw - rw) / 2
    ry = zy + 0.12
    _add_rect(slide, rx, ry, rw, rh, acc, root, font_size=11, text_hex=txt, bold=True)

    nc = len(children)
    cw = min(1.8, (zw - 0.15 * (nc + 1)) / max(nc, 1))
    ch = min(zh * 0.18, 0.65)
    total_cw = cw * nc + 0.12 * (nc - 1)
    csx = zx + (zw - total_cw) / 2
    cy = ry + rh + zh * 0.25

    for i, child in enumerate(children):
        cx = csx + i * (cw + 0.12)
        _add_rect(slide, cx, cy, cw, ch,
                  _darken(acc, 0.2), child, font_size=9, text_hex=txt, bold=False)
        _add_connector(slide, rx + rw / 2, ry + rh, cx + cw / 2, cy, acc)
    return True


# ------------------------------------------------------------------ #
#  Dispatcher                                                         #
# ------------------------------------------------------------------ #

RENDERERS = {
    "process":    _render_process,
    "diagram":    _render_process,
    "cycle":      _render_process,
    "timeline":   _render_timeline,
    "comparison": _render_comparison,
    "matrix":     _render_matrix,
    "orgchart":   _render_orgchart,
}


def render_diagram_in_zone(visual_type: str, prompt: str, slide, style: dict,
                           zone_x: float, zone_y: float,
                           zone_w: float, zone_h: float) -> bool:
    """
    Рисует схему строго в зоне (zone_x, zone_y, zone_w, zone_h) — Inches.
    Возвращает True при успехе, False при ошибке или неизвестном типе.
    """
    fn = RENDERERS.get(visual_type)
    if not fn:
        return False
    try:
        return fn(prompt, slide, style, zone_x, zone_y, zone_w, zone_h)
    except Exception:
        return False


def render_diagram(visual_type: str, prompt: str, slide, style: dict,
                   slide_width=13.33, slide_height=7.5) -> bool:
    """Обратная совместимость — рисует в зоне контента на весь слайд."""
    return render_diagram_in_zone(
        visual_type, prompt, slide, style,
        0.4, 1.7, slide_width - 0.8, slide_height - 2.1,
    )
