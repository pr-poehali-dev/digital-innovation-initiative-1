"""
Экспорт результата генерации в PPTX файл.
Парсит текст с AI, создаёт слайды и возвращает base64-encoded PPTX.
"""
import json
import os
import re
import base64
import io
import boto3
import psycopg2


def get_db():
    conn = psycopg2.connect(os.environ["DATABASE_URL"])
    conn.autocommit = False
    return conn


def get_schema():
    return os.environ.get("MAIN_DB_SCHEMA", "public")


ALLOWED_ORIGINS = {
    "https://raven.moscow",
    "https://www.raven.moscow",
    "https://docmind.ai",
    "https://digital-innovation-initiative-1--preview.poehali.dev",
    "https://poehali.dev",
    "http://localhost:5173",
    "http://localhost:3000",
}


def _is_allowed_origin(origin: str) -> bool:
    """Безопасная проверка origin через urlparse — не raw endswith.
    Защита от попыток типа https://attacker.com/.poehali.dev"""
    if not origin:
        return False
    if origin in ALLOWED_ORIGINS:
        return True
    try:
        from urllib.parse import urlparse
        parsed = urlparse(origin)
        # Только https + допустимый hostname
        if parsed.scheme not in ("https", "http"):
            return False
        hostname = (parsed.hostname or "").lower()
        # Точное совпадение с poehali.dev или его поддоменом
        if hostname == "poehali.dev" or hostname.endswith(".poehali.dev"):
            return True
        return False
    except Exception:
        return False


def cors_headers(origin: str = None):
    """Strict CORS: deny-by-default. Без Allow-Credentials — auth через header, не cookies."""
    headers = {
        "Access-Control-Allow-Methods": "GET, POST, PUT, DELETE, OPTIONS",
        "Access-Control-Allow-Headers": "Content-Type, X-Session-Id",
        "Vary": "Origin",
    }
    if _is_allowed_origin(origin):
        headers["Access-Control-Allow-Origin"] = origin
    return headers


def json_response(data, status=200, origin=None):
    return {
        "statusCode": status,
        "headers": {**cors_headers(origin), "Content-Type": "application/json"},
        "body": json.dumps(data, ensure_ascii=False, default=str),
    }


def get_current_user(conn, session_id):
    if not session_id:
        return None
    schema = get_schema()
    cur = conn.cursor()
    cur.execute(
        f"SELECT u.id, u.email, u.name FROM {schema}.sessions s JOIN {schema}.users u ON u.id = s.user_id WHERE s.id = %s AND s.expires_at > NOW()",
        (session_id,),
    )
    row = cur.fetchone()
    if row:
        return {"id": row[0], "email": row[1], "name": row[2]}
    return None


def parse_slides_from_text(text: str) -> list:
    """Парсит текст AI-результата в список слайдов."""
    slides = []

    # Паттерн 1: "Слайд N: Заголовок" или "Слайд N. Заголовок"
    slide_pattern = re.compile(
        r'(?:Слайд\s+(\d+)[:.]\s*(.+?))\n(.*?)(?=(?:Слайд\s+\d+[:.])|\Z)',
        re.DOTALL | re.IGNORECASE
    )
    matches = list(slide_pattern.finditer(text))

    if matches:
        for m in matches:
            num = int(m.group(1))
            title = m.group(2).strip()
            body = m.group(3).strip()

            # Извлекаем буллеты
            bullets = []
            notes = ""
            for line in body.split('\n'):
                line = line.strip()
                if not line:
                    continue
                if line.lower().startswith(('заметки', 'notes', 'спикер', 'speaker')):
                    notes = re.sub(r'^(заметки|notes|спикер|speaker)[:\s]*', '', line, flags=re.IGNORECASE).strip()
                    continue
                # Убираем маркеры буллетов
                clean = re.sub(r'^[•\-\*\–\—\d+\.]+\s*', '', line).strip()
                if clean:
                    bullets.append(clean)

            slides.append({
                "num": num,
                "title": title,
                "bullets": bullets[:8],
                "notes": notes,
            })
    else:
        # Паттерн 2: заголовки через ## или **Заголовок**
        chunks = re.split(r'\n(?=#{1,3}\s|\*\*[^*]+\*\*\n)', text)
        num = 1
        for chunk in chunks:
            chunk = chunk.strip()
            if not chunk:
                continue
            lines = chunk.split('\n')
            title_raw = lines[0].strip()
            title = re.sub(r'^#{1,3}\s*|\*\*|\*', '', title_raw).strip()
            if not title:
                continue
            bullets = []
            for line in lines[1:]:
                line = line.strip()
                clean = re.sub(r'^[•\-\*\–\—\d+\.]+\s*', '', line).strip()
                clean = re.sub(r'\*\*|\*', '', clean).strip()
                if clean and len(clean) > 3:
                    bullets.append(clean)
            slides.append({
                "num": num,
                "title": title[:80],
                "bullets": bullets[:8],
                "notes": "",
            })
            num += 1

    # Если совсем ничего не распарсилось — делаем один слайд с полным текстом
    if not slides:
        slides = [{
            "num": 1,
            "title": "Результат",
            "bullets": [line.strip() for line in text.split('\n') if line.strip()][:15],
            "notes": "",
        }]

    return slides


# ============================================================
# СТИЛИ PPTX — пресеты и извлечение из образца
# ============================================================

STYLE_PRESETS = {
    "dark_corporate": {
        "bg": (0x1a, 0x20, 0x35), "title": (0xFF, 0xFF, 0xFF),
        "bullet": (0xE2, 0xE8, 0xF0), "accent": (0x60, 0x8B, 0xC4),
        "num": (0x94, 0xA3, 0xB8), "font": "Calibri",
    },
    "light_minimal": {
        "bg": (0xFF, 0xFF, 0xFF), "title": (0x1F, 0x29, 0x37),
        "bullet": (0x37, 0x41, 0x51), "accent": (0x25, 0x63, 0xEB),
        "num": (0x9C, 0xA3, 0xAF), "font": "Inter",
    },
    "academic": {
        "bg": (0xFA, 0xF7, 0xF2), "title": (0x1F, 0x2D, 0x3D),
        "bullet": (0x33, 0x40, 0x55), "accent": (0x8B, 0x4D, 0x2F),
        "num": (0x8E, 0x8E, 0x8E), "font": "Times New Roman",
    },
    "marketing": {
        "bg": (0xFD, 0xF2, 0xF8), "title": (0x86, 0x19, 0x4E),
        "bullet": (0x4A, 0x14, 0x4A), "accent": (0xEC, 0x48, 0x99),
        "num": (0xC0, 0x84, 0x9C), "font": "Montserrat",
    },
    "scientific": {
        "bg": (0xF0, 0xF9, 0xFF), "title": (0x07, 0x2C, 0x4D),
        "bullet": (0x1E, 0x3A, 0x5F), "accent": (0x05, 0x96, 0x69),
        "num": (0x6B, 0x72, 0x80), "font": "Roboto",
    },
}


def _rgb_to_tuple(rgb):
    """Безопасно конвертирует RGBColor / hex-string в кортеж (r,g,b)."""
    try:
        s = str(rgb)
        if len(s) == 6:
            return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except BaseException:
        pass
    return None


def extract_style_from_template(template_bytes: bytes) -> dict:
    """Извлекает цветовую палитру и шрифт из PPTX-образца (роль template).
    Любая ошибка — fallback на dark_corporate."""
    fallback = STYLE_PRESETS["dark_corporate"]
    if not template_bytes:
        return fallback

    try:
        from pptx import Presentation as PptxPresentation
        from io import BytesIO as _BIO
        prs = PptxPresentation(_BIO(template_bytes))
    except BaseException:
        return fallback

    bg_color = None
    title_color = None
    body_color = None
    main_font = None

    try:
        slides_iter = list(prs.slides)[:5]
    except BaseException:
        slides_iter = []

    for slide in slides_iter:
        # Фон слайда
        try:
            fill = slide.background.fill
            if bg_color is None:
                rgb_t = _rgb_to_tuple(fill.fore_color.rgb)
                if rgb_t:
                    bg_color = rgb_t
        except BaseException:
            pass

        # Текстовые рамки
        try:
            shapes = list(slide.shapes)
        except BaseException:
            shapes = []

        for shape in shapes:
            try:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
            except BaseException:
                continue

            try:
                paragraphs = list(tf.paragraphs)
            except BaseException:
                paragraphs = []

            for para in paragraphs:
                try:
                    runs = list(para.runs)
                except BaseException:
                    runs = []
                for run in runs:
                    try:
                        if not run.text or not run.text.strip():
                            continue
                    except BaseException:
                        continue
                    try:
                        font = run.font
                    except BaseException:
                        continue
                    # Шрифт
                    try:
                        if main_font is None and font.name:
                            main_font = font.name
                    except BaseException:
                        pass
                    # Цвет
                    try:
                        rgb_t = _rgb_to_tuple(font.color.rgb)
                        if rgb_t:
                            try:
                                size_pt = font.size.pt if font.size else 18
                            except BaseException:
                                size_pt = 18
                            if size_pt >= 24 and title_color is None:
                                title_color = rgb_t
                            elif size_pt < 24 and body_color is None:
                                body_color = rgb_t
                    except BaseException:
                        pass

    return {
        "bg": bg_color or fallback["bg"],
        "title": title_color or fallback["title"],
        "bullet": body_color or fallback["bullet"],
        "accent": title_color or fallback["accent"],
        "num": fallback["num"],
        "font": main_font or fallback["font"],
    }


def resolve_style(style_preset: str, template_bytes: bytes = None) -> dict:
    """По имени пресета или из шаблона — возвращает палитру."""
    if style_preset == "from_template" and template_bytes:
        return extract_style_from_template(template_bytes)
    if style_preset in STYLE_PRESETS:
        return STYLE_PRESETS[style_preset]
    return STYLE_PRESETS["dark_corporate"]


def _choose_layout_mode(visual_type: str) -> str:
    """Выбирает layout-режим по типу визуала."""
    if visual_type in ("timeline", "cycle"):
        return "title_text_top_timeline_bottom"
    if visual_type in ("diagram", "process", "comparison", "matrix", "orgchart", "image"):
        return "title_text_left_visual_right"
    return "title_text_left_visual_right"


def _draw_slide_chrome(slide, slide_num: int, title_text: str, style: dict):
    """Рисует общие элементы слайда: фон, полоска, номер, заголовок, разделитель."""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    COLOR_BG = RGBColor(*style["bg"])
    COLOR_TITLE = RGBColor(*style["title"])
    COLOR_ACCENT = RGBColor(*style["accent"])
    COLOR_NUM = RGBColor(*style["num"])
    FONT_NAME = style.get("font", "Calibri")

    # Фон
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_BG

    # Акцентная полоска слева
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    # Номер слайда
    num_box = slide.shapes.add_textbox(Inches(12.3), Inches(0.2), Inches(0.8), Inches(0.5))
    tf = num_box.text_frame
    tf.text = str(slide_num)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.color.rgb = COLOR_NUM
    run.font.bold = False
    run.font.name = FONT_NAME

    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(1.1))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.text = title_text
    p = tf.paragraphs[0]
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = COLOR_TITLE
    run.font.name = FONT_NAME

    # Разделитель
    div = slide.shapes.add_shape(1, Inches(0.4), Inches(1.35), Inches(12.5), Emu(32000))
    div.fill.solid()
    div.fill.fore_color.rgb = COLOR_ACCENT
    div.line.fill.background()


def _draw_bullets(slide, bullets: list, x: float, y: float, w: float, h: float,
                  max_bullets: int, style: dict, font_size: int = 16):
    """Рисует буллеты в заданной зоне. Обрезает до max_bullets, уменьшает шрифт если надо."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if not bullets:
        return

    COLOR_BULLET = RGBColor(*style["bullet"])
    COLOR_ACCENT = RGBColor(*style["accent"])
    FONT_NAME = style.get("font", "Calibri")

    # Ограничиваем количество и длину
    capped = bullets[:max_bullets]
    # Если буллет длиннее 120 символов — обрезаем
    capped = [b[:120] + ("…" if len(b) > 120 else "") for b in capped]

    content_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(capped):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(3)

        run_bullet = p.add_run()
        run_bullet.text = "▸  "
        run_bullet.font.size = Pt(max(10, font_size - 4))
        run_bullet.font.color.rgb = COLOR_ACCENT
        run_bullet.font.bold = True
        run_bullet.font.name = FONT_NAME

        run_text = p.add_run()
        run_text.text = bullet
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = COLOR_BULLET
        run_text.font.bold = False
        run_text.font.name = FONT_NAME


def build_pptx(slides: list, title: str, style: dict = None, visual_plan: list = None) -> bytes:
    """
    Строит PPTX. Если передан visual_plan — применяет layout-движок:
      - title_text_left_visual_right: текст слева 42%, визуал справа 52%
      - title_text_top_timeline_bottom: текст сверху 30%, визуал внизу 65%
    Визуалы рисуются сразу без overlap с текстом.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if style is None:
        style = STYLE_PRESETS["dark_corporate"]

    COLOR_BG = RGBColor(*style["bg"])
    COLOR_TITLE = RGBColor(*style["title"])
    COLOR_BULLET = RGBColor(*style["bullet"])
    COLOR_ACCENT = RGBColor(*style["accent"])
    COLOR_NUM = RGBColor(*style["num"])
    FONT_NAME = style.get("font", "Calibri")

    # Строим индекс: slide_index → visual_plan_item (1-based, title=0)
    visual_index: dict = {}
    if visual_plan:
        for vp in visual_plan:
            si = int(vp.get("slide_index", 0))
            if si > 0:
                visual_index[si] = vp

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # slide_data["num"] начинается с 1 — совпадает с visual_plan slide_index
    for slide_data in slides:
        slide = prs.slides.add_slide(blank_layout)
        snum = slide_data["num"]

        # --- Chrome (фон, полоска, номер, заголовок, разделитель) ---
        _draw_slide_chrome(slide, snum, slide_data["title"], style)

        bullets = slide_data.get("bullets") or []
        vp = visual_index.get(snum)

        if not vp:
            # ===== Обычный слайд — полная ширина =====
            _draw_bullets(slide, bullets,
                          x=0.55, y=1.55, w=12.3, h=5.6,
                          max_bullets=8, style=style, font_size=17)
        else:
            # ===== Слайд с визуалом — layout-движок =====
            visual_type = vp.get("visual_type", "diagram")
            render_mode = vp.get("render_mode", "pptx_shapes")
            prompt = vp.get("source_prompt", "")
            layout_mode = _choose_layout_mode(visual_type)
            vp["layout_mode"] = layout_mode   # сохраняем в metadata

            # ------ LAYOUT A: title_text_left_visual_right ------
            if layout_mode == "title_text_left_visual_right":
                # Зоны: текст x=0.4..5.7 (5.3"), визуал x=5.9..13.1 (7.2")
                TEXT_X, TEXT_Y, TEXT_W, TEXT_H = 0.4, 1.55, 5.3, 5.6
                VIS_X,  VIS_Y,  VIS_W,  VIS_H  = 5.9, 1.55, 7.0, 5.6

                _draw_bullets(slide, bullets,
                              x=TEXT_X, y=TEXT_Y, w=TEXT_W, h=TEXT_H,
                              max_bullets=5, style=style, font_size=15)

                # Рисуем визуал в правой зоне
                _render_visual_in_zone(slide, vp, style,
                                       VIS_X, VIS_Y, VIS_W, VIS_H)

            # ------ LAYOUT B: title_text_top_timeline_bottom ------
            elif layout_mode == "title_text_top_timeline_bottom":
                # Текст в верхней полосе; timeline занимает нижние 65%
                TEXT_X, TEXT_Y, TEXT_W, TEXT_H = 0.4, 1.55, 12.5, 1.6
                VIS_X,  VIS_Y,  VIS_W,  VIS_H  = 0.4, 3.25, 12.5, 4.0

                _draw_bullets(slide, bullets,
                              x=TEXT_X, y=TEXT_Y, w=TEXT_W, h=TEXT_H,
                              max_bullets=3, style=style, font_size=14)

                _render_visual_in_zone(slide, vp, style,
                                       VIS_X, VIS_Y, VIS_W, VIS_H)

        # Заметки спикера
        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    # Титульный слайд в начало
    title_slide = prs.slides.add_slide(blank_layout)

    bg = title_slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_BG

    # Большой акцент
    bar = title_slide.shapes.add_shape(1, Inches(0), Inches(3.2), Inches(13.33), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()

    # Заголовок презентации
    tbox = title_slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11), Inches(1.5))
    tf = tbox.text_frame
    tf.word_wrap = True
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = COLOR_TITLE
    run.font.name = FONT_NAME

    # Подзаголовок DocMind AI
    sbox = title_slide.shapes.add_textbox(Inches(1), Inches(3.6), Inches(11), Inches(0.8))
    tf = sbox.text_frame
    tf.text = "DocMind AI"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(18)
    run.font.color.rgb = COLOR_ACCENT
    run.font.bold = False
    run.font.name = FONT_NAME

    # Перемещаем титульный слайд в начало
    xml_slides = prs.slides._sldIdLst
    xml_slides.insert(0, xml_slides[-1])
    xml_slides[-1].getparent().remove(xml_slides[-1])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _load_image_from_s3(s3_key: str):
    """Скачивает файл из S3. Возвращает bytes или None при ошибке."""
    try:
        s3_client = boto3.client(
            "s3",
            endpoint_url="https://bucket.poehali.dev",
            aws_access_key_id=os.environ["AWS_ACCESS_KEY_ID"],
            aws_secret_access_key=os.environ["AWS_SECRET_ACCESS_KEY"],
        )
        obj = s3_client.get_object(Bucket="files", Key=s3_key)
        return obj["Body"].read()
    except Exception:
        return None


def _render_visual_in_zone(slide, vp: dict, style: dict,
                           zone_x: float, zone_y: float, zone_w: float, zone_h: float):
    """
    Рисует визуал в заданной зоне (Inches).

    Приоритет active_asset:
      1. user_override_s3_key  (user_image / user_override)
      2. active_asset_s3_key
      3. asset_s3_key          (ai_image)
      4. pptx_shapes           (схема через diagram_renderer)
      5. fallback placeholder
    """
    from pptx.util import Inches as _In

    visual_type = vp.get("visual_type", "diagram")
    render_mode = vp.get("render_mode", "pptx_shapes")
    prompt = vp.get("source_prompt", "")

    # --- Определяем S3-ключ для картинки: user override > active > ai ---
    image_s3_key = (
        vp.get("user_override_s3_key")
        or vp.get("active_asset_s3_key")
        or (vp.get("asset_s3_key") if render_mode == "ai_image" else None)
    )

    # Если есть user-override или ai-картинка — вставляем как изображение
    if image_s3_key and render_mode in ("user_image", "ai_image", "user_override"):
        img_bytes = _load_image_from_s3(image_s3_key)
        if img_bytes:
            try:
                slide.shapes.add_picture(
                    io.BytesIO(img_bytes),
                    _In(zone_x), _In(zone_y), _In(zone_w), _In(zone_h),
                )
                vp["generation_status"] = "rendered"
                return
            except Exception:
                pass
        # Картинка есть в плане, но не загрузилась
        _add_image_placeholder(slide, zone_x, zone_y, zone_w, zone_h, prompt, style)
        vp["generation_status"] = "image_pending"
        return

    # --- Схемы через diagram_renderer ---
    if render_mode == "pptx_shapes":
        try:
            from diagram_renderer import render_diagram_in_zone
            ok = render_diagram_in_zone(
                visual_type, prompt, slide, style,
                zone_x, zone_y, zone_w, zone_h,
            )
            if ok:
                vp["generation_status"] = "rendered"
                return
        except Exception:
            pass
        _add_visual_fallback(slide, zone_x, zone_y, zone_w, zone_h, visual_type, prompt, style)
        vp["generation_status"] = "fallback"
        return

    # --- ai_image без s3_key ---
    if render_mode == "ai_image":
        _add_image_placeholder(slide, zone_x, zone_y, zone_w, zone_h, prompt, style)
        vp["generation_status"] = "image_pending"
        return

    _add_visual_fallback(slide, zone_x, zone_y, zone_w, zone_h, visual_type, prompt, style)
    vp["generation_status"] = "fallback"


def _add_image_placeholder(slide, x, y, w, h, prompt, style):
    """Placeholder для ещё не сгенерированной картинки."""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    accent = style.get("accent", (0x60, 0x8B, 0xC4))
    title_rgb = style.get("title", (0xFF, 0xFF, 0xFF))

    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    r, g, b = accent
    box.fill.fore_color.rgb = RGBColor(max(0, r - 60), max(0, g - 40), max(0, b - 20))
    box.line.color.rgb = RGBColor(*accent)
    box.line.width = Emu(20000)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"🖼 Картинка\n{prompt[:90]}"
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(*title_rgb)
    run.font.italic = True


def apply_visual_plan(prs, visual_plan: list, style: dict):
    """
    Теперь build_pptx сам встраивает визуалы через layout-движок.
    apply_visual_plan — резервный проход: ничего не делает, visual_plan уже обработан.
    Оставлен для обратной совместимости.
    """
    pass


def _add_visual_fallback(slide, x, y, w, h, visual_type, prompt, style):
    """Аккуратный fallback — рамка с подписью типа и промптом."""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    accent = style.get("accent", (0x60, 0x8B, 0xC4))
    title_rgb = style.get("title", (0xFF, 0xFF, 0xFF))

    # Рамка-placeholder
    box = slide.shapes.add_shape(
        1,  # RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*accent)
    box.fill.fore_color.rgb = RGBColor(accent[0] // 4, accent[1] // 4 + 20, accent[2] // 4 + 30)
    box.line.color.rgb = RGBColor(*accent)
    box.line.width = Emu(18000)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[{visual_type.upper()}]\n{prompt[:80]}"
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(*title_rgb)
    run.font.italic = True


def handler(event: dict, context) -> dict:
    origin = (event.get("headers") or {}).get("Origin") or (event.get("headers") or {}).get("origin")

    if event.get("httpMethod") == "OPTIONS":
        return {"statusCode": 200, "headers": cors_headers(origin), "body": ""}

    method = event.get("httpMethod", "GET")
    body = {}
    if event.get("body"):
        try:
            body = json.loads(event["body"])
        except Exception:
            pass

    session_id = event.get("headers", {}).get("X-Session-Id", "")
    conn = get_db()
    schema = get_schema()

    try:
        user = get_current_user(conn, session_id)
        if not user:
            return json_response({"error": "Не авторизован"}, 401, origin=origin)

        if method != "POST":
            return json_response({"error": "Method not allowed"}, 405, origin=origin)

        run_id = body.get("run_id")
        if not run_id:
            return json_response({"error": "Нужен run_id"}, 400, origin=origin)

        cur = conn.cursor()

        # Загружаем run
        cur.execute(
            f"""SELECT gr.result_json, gr.task_id
                FROM {schema}.generation_runs gr
                WHERE gr.id = %s""",
            (run_id,),
        )
        row = cur.fetchone()
        if not row:
            return json_response({"error": "Результат не найден"}, 404, origin=origin)

        result_json, task_id = row

        # Загружаем задание для заголовка + style_preset
        cur.execute(
            f"SELECT title, topic, project_id, style_preset FROM {schema}.tasks WHERE id = %s",
            (task_id,),
        )
        task_row = cur.fetchone()
        if not task_row:
            return json_response({"error": "Задание не найдено"}, 404, origin=origin)

        task_title, task_topic, project_id, style_preset = task_row

        # Проверяем доступ
        cur.execute(
            f"SELECT role FROM {schema}.project_members WHERE project_id = %s AND user_id = %s",
            (project_id, user["id"]),
        )
        if not cur.fetchone():
            return json_response({"error": "Нет доступа"}, 403, origin=origin)

        # Извлекаем текст
        content = ""
        if result_json:
            try:
                content = json.loads(result_json).get("content", "")
            except Exception:
                content = result_json

        if not content:
            return json_response({"error": "Нет контента для экспорта"}, 400, origin=origin)

        # Парсим
        slides = parse_slides_from_text(content)
        pptx_title = task_topic or task_title or "Презентация"

        # Определяем стиль: пресет или из образца (роль template)
        template_bytes = None
        effective_preset = style_preset or "from_template"

        if effective_preset == "from_template":
            # Ищем документ с ролью template в задании
            cur.execute(
                f"""SELECT d.s3_key, d.file_type, d.original_name
                    FROM {schema}.task_documents td
                    JOIN {schema}.documents d ON d.id = td.document_id
                    WHERE td.task_id = %s AND td.role = 'template' AND d.archived_at IS NULL
                    ORDER BY td.priority DESC
                    LIMIT 1""",
                (task_id,),
            )
            tpl_row = cur.fetchone()
            if tpl_row and tpl_row[0] and "pptx" in (tpl_row[1] or "").lower():
                try:
                    s3 = boto3.client(
                        "s3",
                        endpoint_url="https://bucket.poehali.dev",
                        aws_access_key_id=os.environ["AWS_ACCESS_KEY_ID"],
                        aws_secret_access_key=os.environ["AWS_SECRET_ACCESS_KEY"],
                    )
                    obj = s3.get_object(Bucket="files", Key=tpl_row[0])
                    template_bytes = obj["Body"].read()
                except Exception:
                    template_bytes = None
            # Если не нашли шаблон — fallback на dark_corporate
            if not template_bytes:
                effective_preset = "dark_corporate"

        style = resolve_style(effective_preset, template_bytes)

        # visual_plan из result_json
        visual_plan = []
        if result_json:
            try:
                rj = json.loads(result_json)
                visual_plan = rj.get("visual_plan") or []
            except Exception:
                pass

        # build_pptx теперь принимает visual_plan и строит layout сразу
        pptx_bytes = build_pptx(slides, pptx_title, style=style, visual_plan=visual_plan)

        # Возвращаем base64
        pptx_b64 = base64.b64encode(pptx_bytes).decode("utf-8")
        filename = f"{pptx_title[:40].replace(' ', '_')}.pptx"

        return json_response({
            "filename": filename,
            "file_data": pptx_b64,
            "slides_count": len(slides),
        }, origin=origin)

    finally:
        conn.close()