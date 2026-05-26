"""
Visual Renderer — генерирует визуальные элементы для PPTX.

A. Картинки (visual_type=image)     → YandexArt → PNG → S3
B. Схемы (diagram/timeline/process/...) → python-pptx shapes → bytes patch

Все функции возвращают:
  {"ok": True, "s3_key": ..., "asset_bytes": bytes | None, "warning": None | str}
или
  {"ok": False, "warning": str}
"""
import os
import io
import re
import time
import json
import base64
import urllib.request
import urllib.error
import boto3


def get_s3():
    return boto3.client(
        "s3",
        endpoint_url="https://bucket.poehali.dev",
        aws_access_key_id=os.environ["AWS_ACCESS_KEY_ID"],
        aws_secret_access_key=os.environ["AWS_SECRET_ACCESS_KEY"],
    )


# ============================================================
# A. YandexArt — генерация картинки
# ============================================================

def _yandex_art_request(prompt: str, style_hint: str = "") -> dict:
    """Синхронный (polling) вызов YandexArt API."""
    api_key = os.environ.get("YANDEX_GPT_API_KEY", "")
    folder_id = os.environ.get("YANDEX_FOLDER_ID", "")
    if not api_key or not folder_id:
        return {"ok": False, "warning": "YandexArt API key не настроен"}

    # Собираем промпт — clean illustration без текста внутри
    full_prompt = prompt.strip()
    if style_hint:
        full_prompt = f"{full_prompt}, {style_hint}"
    full_prompt += ", professional illustration, no text inside, clean design, high quality"

    payload = json.dumps({
        "modelUri": f"art://{folder_id}/yandex-art/latest",
        "generationOptions": {"seed": 42, "aspectRatio": {"widthRatio": "16", "heightRatio": "9"}},
        "messages": [{"weight": "1", "text": full_prompt}],
    }).encode()

    req = urllib.request.Request(
        "https://llm.api.cloud.yandex.net/foundationModels/v1/imageGenerationAsync",
        data=payload,
        headers={"Authorization": f"Api-Key {api_key}", "Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=30) as resp:
            result = json.loads(resp.read())
        operation_id = result.get("id")
        if not operation_id:
            return {"ok": False, "warning": f"YandexArt: нет operation_id. Ответ: {str(result)[:200]}"}
    except Exception as e:
        return {"ok": False, "warning": f"YandexArt запрос: {e}"}

    # Polling до 60 секунд
    poll_url = f"https://llm.api.cloud.yandex.net/operations/{operation_id}"
    poll_req = urllib.request.Request(
        poll_url,
        headers={"Authorization": f"Api-Key {api_key}"},
        method="GET",
    )
    for _ in range(20):
        time.sleep(3)
        try:
            with urllib.request.urlopen(poll_req, timeout=15) as resp:
                status = json.loads(resp.read())
            if status.get("done"):
                img_b64 = status.get("response", {}).get("image", "")
                if img_b64:
                    return {"ok": True, "image_bytes": base64.b64decode(img_b64)}
                return {"ok": False, "warning": "YandexArt: done=True но нет image в ответе"}
            if status.get("error"):
                return {"ok": False, "warning": f"YandexArt error: {status['error'].get('message', str(status['error'])[:200])}"}
        except Exception as e:
            return {"ok": False, "warning": f"YandexArt polling: {e}"}

    return {"ok": False, "warning": "YandexArt: timeout — картинка не сгенерировалась за 60 сек"}


def render_image(prompt: str, s3_key: str, style_hint: str = "") -> dict:
    """Генерирует PNG через YandexArt и сохраняет в S3."""
    result = _yandex_art_request(prompt, style_hint)
    if not result.get("ok"):
        return result

    img_bytes = result["image_bytes"]
    try:
        s3 = get_s3()
        s3.put_object(Bucket="files", Key=s3_key, Body=img_bytes, ContentType="image/png")
    except Exception as e:
        return {"ok": False, "warning": f"S3 upload: {e}"}

    aws_key = os.environ.get("AWS_ACCESS_KEY_ID", "")
    cdn_url = f"https://cdn.poehali.dev/projects/{aws_key}/bucket/{s3_key}"
    return {"ok": True, "s3_key": s3_key, "asset_url": cdn_url, "image_bytes": img_bytes}


# ============================================================
# B. Схемы через python-pptx shapes
# ============================================================

def _rgb(hex_str: str):
    from pptx.dml.color import RGBColor
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _parse_steps_from_prompt(prompt: str) -> list:
    """
    Извлекает список шагов из промпта.
    Поддерживает: → / -> / em-dash — / нумерацию / запятые.
    Автоматически убирает prefix "N этапов:".
    """
    norm = prompt.replace("→", "->").replace("—", ",")

    if "->" in norm:
        parts = norm.split("->")
        cleaned = []
        for s in parts:
            s = s.strip().strip(".,;")
            if ":" in s:
                prefix, _, rest = s.partition(":")
                if len(prefix.strip()) < 40:
                    s = rest.strip()
            if s:
                cleaned.append(s[:60])
        if len(cleaned) >= 2:
            return cleaned[:8]

    numbered = re.findall(r'(?:^|\b)(\d+)[.\)]\s*([^\d\n;]+?)(?=\s*\d+[.\)]|\Z)', norm)
    if len(numbered) >= 2:
        return [v.strip().strip(",")[:60] for _, v in numbered][:8]

    parts = [p.strip()[:60] for p in norm.split(",") if p.strip()]
    if len(parts) >= 2:
        return parts[:8]

    return [prompt[:60]]


def _add_rounded_rect(slide, x, y, w, h, fill_hex, text, font_size=11, text_color="#FFFFFF", bold=False):
    """Добавляет прямоугольник со скруглёнными углами и текстом."""
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.color.rgb = _rgb(fill_hex)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(text_color)
    return shape


def _add_arrow(slide, x1, y1, x2, y2, color_hex="#94A3B8"):
    """Рисует линию-стрелку между двумя точками через add_connector."""
    from pptx.util import Inches, Emu
    # MSO_CONNECTOR_TYPE.STRAIGHT = 1
    connector = slide.shapes.add_connector(
        1,  # STRAIGHT connector
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    connector.line.color.rgb = _rgb(color_hex)
    connector.line.width = Emu(18000)  # 1.5pt
    return connector


def render_process_diagram(prompt: str, slide, style: dict, slide_width=13.33, slide_height=7.5):
    """Рисует процессную схему (горизонтальные блоки со стрелками)."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    steps = _parse_steps_from_prompt(prompt)
    n = len(steps)
    if n == 0:
        return False

    accent = style.get("accent", (0x60, 0x8B, 0xC4))
    accent_hex = "#{:02X}{:02X}{:02X}".format(*accent)
    bg = style.get("bg", (0x1a, 0x20, 0x35))
    bg_hex = "#{:02X}{:02X}{:02X}".format(*bg)
    title_hex = "#{:02X}{:02X}{:02X}".format(*style.get("title", (0xFF, 0xFF, 0xFF)))

    # Параметры сетки
    margin_x = 0.6
    area_y = 2.2
    area_h = 3.5
    total_w = slide_width - margin_x * 2
    block_w = min(2.2, (total_w - 0.3 * (n - 1)) / n)
    gap = (total_w - block_w * n) / max(n - 1, 1) if n > 1 else 0
    block_h = 1.5
    block_y = area_y + (area_h - block_h) / 2

    for i, step in enumerate(steps):
        x = margin_x + i * (block_w + gap)
        # Чередуем оттенки
        fill = accent_hex if i % 2 == 0 else _darken(accent_hex, 0.2)
        _add_rounded_rect(slide, x, block_y, block_w, block_h,
                          fill, step, font_size=10, text_color=title_hex, bold=True)

        # Стрелка вправо (кроме последнего)
        if i < n - 1:
            ax = x + block_w + 0.02
            ay = block_y + block_h / 2
            ax2 = ax + gap - 0.04
            _add_arrow(slide, ax, ay, ax2, ay, color_hex=accent_hex)

    return True


def render_timeline_diagram(prompt: str, slide, style: dict, slide_width=13.33, slide_height=7.5):
    """Рисует таймлайн — горизонтальная ось с точками-этапами."""
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    steps = _parse_steps_from_prompt(prompt)
    n = len(steps)
    if n == 0:
        return False

    accent = style.get("accent", (0x60, 0x8B, 0xC4))
    accent_hex = "#{:02X}{:02X}{:02X}".format(*accent)
    title_hex = "#{:02X}{:02X}{:02X}".format(*style.get("title", (0xFF, 0xFF, 0xFF)))

    margin_x = 0.8
    axis_y = 4.0
    total_w = slide_width - margin_x * 2

    # Горизонтальная ось
    axis_line = slide.shapes.add_connector(1, Inches(margin_x), Inches(axis_y),
                                            Inches(slide_width - margin_x), Inches(axis_y))
    axis_line.line.color.rgb = _rgb(accent_hex)
    axis_line.line.width = Emu(28000)

    step_gap = total_w / max(n - 1, 1) if n > 1 else total_w / 2

    for i, step in enumerate(steps):
        sx = margin_x + i * step_gap

        # Точка на оси (OVAL)
        from pptx.enum.shapes import MSO_SHAPE as _MS
        dot = slide.shapes.add_shape(_MS.OVAL, Inches(sx - 0.12), Inches(axis_y - 0.12),
                                     Inches(0.24), Inches(0.24))
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(accent_hex)
        dot.line.fill.background()

        # Метка чередуется вверх/вниз
        if i % 2 == 0:
            ty = axis_y - 1.4
        else:
            ty = axis_y + 0.3

        label = slide.shapes.add_textbox(Inches(sx - 0.9), Inches(ty), Inches(1.8), Inches(0.9))
        tf = label.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step[:50]
        run.font.size = Pt(9)
        run.font.color.rgb = _rgb(title_hex)
        run.font.bold = (i == 0 or i == n - 1)

        # Линия от точки к метке
        ly1 = axis_y - 0.12 if i % 2 == 0 else axis_y + 0.12
        ly2 = ty + (0.9 if i % 2 == 0 else 0)
        _add_arrow(slide, sx, ly1, sx, ly2, color_hex=accent_hex)

    return True


def render_comparison_diagram(prompt: str, slide, style: dict, slide_width=13.33, slide_height=7.5):
    """Два столбца для сравнения."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    accent = style.get("accent", (0x60, 0x8B, 0xC4))
    accent_hex = "#{:02X}{:02X}{:02X}".format(*accent)
    title_hex = "#{:02X}{:02X}{:02X}".format(*style.get("title", (0xFF, 0xFF, 0xFF)))

    # Пробуем выделить два объекта сравнения из промпта
    parts = re.split(r'\bи\b|vs\.?|versus|\bпротив\b', prompt, flags=re.IGNORECASE)
    if len(parts) >= 2:
        left_label = parts[0].strip()[:40]
        right_label = parts[1].strip()[:40]
    else:
        left_label = "Вариант А"
        right_label = "Вариант Б"

    col_w = 5.0
    col_h = 4.0
    col_y = 2.3
    left_x = 0.8
    right_x = slide_width - 0.8 - col_w

    _add_rounded_rect(slide, left_x, col_y, col_w, col_h,
                      accent_hex, left_label, font_size=13, text_color=title_hex, bold=True)
    _add_rounded_rect(slide, right_x, col_y, col_w, col_h,
                      _darken(accent_hex, 0.25), right_label, font_size=13, text_color=title_hex, bold=True)

    # VS посередине
    vs_box = slide.shapes.add_textbox(
        Inches(slide_width / 2 - 0.35), Inches(col_y + col_h / 2 - 0.3), Inches(0.7), Inches(0.6)
    )
    tf = vs_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "VS"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = _rgb(title_hex)

    return True


def render_matrix_diagram(prompt: str, slide, style: dict, slide_width=13.33, slide_height=7.5):
    """2×2 матрица (риски, приоритеты и т.д.)."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    accent = style.get("accent", (0x60, 0x8B, 0xC4))
    accent_hex = "#{:02X}{:02X}{:02X}".format(*accent)
    title_hex = "#{:02X}{:02X}{:02X}".format(*style.get("title", (0xFF, 0xFF, 0xFF)))

    # Цвета квадрантов
    q_colors = [
        _darken(accent_hex, 0.3),  # верх-лево
        accent_hex,                 # верх-право
        _lighten(accent_hex, 0.3), # низ-лево
        _darken(accent_hex, 0.1),  # низ-право
    ]
    labels = ["Низкий / Низкий", "Высокий / Низкий",
              "Низкий / Высокий", "Высокий / Высокий"]

    size = 2.3
    start_x = (slide_width - size * 2 - 0.1) / 2
    start_y = 2.2

    for idx, (color, label) in enumerate(zip(q_colors, labels)):
        col = idx % 2
        row = idx // 2
        x = start_x + col * (size + 0.1)
        y = start_y + row * (size + 0.1)
        _add_rounded_rect(slide, x, y, size, size, color, label,
                          font_size=10, text_color=title_hex)

    return True


def render_orgchart_diagram(prompt: str, slide, style: dict, slide_width=13.33, slide_height=7.5):
    """Простая оргсхема: один корень → ветки."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    accent = style.get("accent", (0x60, 0x8B, 0xC4))
    accent_hex = "#{:02X}{:02X}{:02X}".format(*accent)
    title_hex = "#{:02X}{:02X}{:02X}".format(*style.get("title", (0xFF, 0xFF, 0xFF)))

    steps = _parse_steps_from_prompt(prompt)
    if not steps:
        return False

    root = steps[0]
    children = steps[1:] if len(steps) > 1 else ["Команда"]

    root_w = 3.0
    root_h = 0.9
    root_x = (slide_width - root_w) / 2
    root_y = 2.0
    _add_rounded_rect(slide, root_x, root_y, root_w, root_h,
                      accent_hex, root, font_size=12, text_color=title_hex, bold=True)

    # Дочерние узлы
    child_w = min(2.4, (slide_width - 1.6) / len(children))
    total_cw = child_w * len(children) + 0.2 * (len(children) - 1)
    child_start_x = (slide_width - total_cw) / 2
    child_y = root_y + root_h + 0.9

    for i, child in enumerate(children):
        cx = child_start_x + i * (child_w + 0.2)
        _add_rounded_rect(slide, cx, child_y, child_w, 0.8,
                          _darken(accent_hex, 0.2), child, font_size=9, text_color=title_hex)
        # Линия от корня к ребёнку
        _add_arrow(slide, root_x + root_w / 2, root_y + root_h,
                   cx + child_w / 2, child_y, color_hex=accent_hex)

    return True


def render_diagram(visual_type: str, prompt: str, slide, style: dict,
                   slide_width=13.33, slide_height=7.5) -> bool:
    """
    Диспетчер схем. Возвращает True если успешно нарисовано.
    """
    renderers = {
        "process":    render_process_diagram,
        "diagram":    render_process_diagram,
        "timeline":   render_timeline_diagram,
        "comparison": render_comparison_diagram,
        "matrix":     render_matrix_diagram,
        "orgchart":   render_orgchart_diagram,
        "cycle":      render_process_diagram,   # cycle как process
    }
    fn = renderers.get(visual_type)
    if fn:
        try:
            return fn(prompt, slide, style, slide_width, slide_height)
        except Exception as e:
            return False
    return False


# ------------------------------------------------------------------ #
#  Утилиты цвета                                                      #
# ------------------------------------------------------------------ #

def _darken(hex_color: str, factor: float) -> str:
    """Затемняет цвет на factor (0..1)."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = max(0, int(r * (1 - factor)))
    g = max(0, int(g * (1 - factor)))
    b = max(0, int(b * (1 - factor)))
    return f"#{r:02X}{g:02X}{b:02X}"


def _lighten(hex_color: str, factor: float) -> str:
    """Осветляет цвет на factor."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = min(255, int(r + (255 - r) * factor))
    g = min(255, int(g + (255 - g) * factor))
    b = min(255, int(b + (255 - b) * factor))
    return f"#{r:02X}{g:02X}{b:02X}"