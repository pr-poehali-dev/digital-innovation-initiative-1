"""
Presentation Audit v2 — проверка готовой презентации на соответствие документам.

Действия:
  audit.run     — запустить анализ PPTX против документов с ролями
  audit.get     — получить результат аудита по run_id
  audit.list    — список аудитов проекта
"""
import json
import os
import io
import uuid
import logging
import traceback
import psycopg2
from botocore.exceptions import ClientError

logging.basicConfig(level=logging.INFO)
log = logging.getLogger("audit")

ROLE_PRIORITY = {
    "standard":  1,
    "criteria":  2,
    "source":    3,
    "material":  4,
    "template":  5,
    "example":   6,
}

CORS_HEADERS = {
    "Access-Control-Allow-Origin": "*",
    "Access-Control-Allow-Methods": "GET, POST, PUT, DELETE, OPTIONS",
    "Access-Control-Allow-Headers": "Content-Type, Authorization, X-Session-Id, X-User-Id, X-Auth-Token",
    "Access-Control-Max-Age": "86400",
}


def get_schema() -> str:
    val = os.environ.get("MAIN_DB_SCHEMA", "").strip()
    if val:
        return val
    # Fallback: жёстко зашитая схема проекта
    return "t_p61016064_digital_innovation_i"


def get_db():
    conn = psycopg2.connect(os.environ["DATABASE_URL"])
    conn.autocommit = False
    return conn


def ok_resp(data, origin=None):
    return {
        "statusCode": 200,
        "headers": {**CORS_HEADERS, "Content-Type": "application/json"},
        "body": json.dumps({"ok": True, "data": data}, ensure_ascii=False, default=str),
    }


def err_resp(msg, status=400, origin=None, detail: dict = None):
    body = {"ok": False, "error": msg}
    if detail:
        body["detail"] = detail
    return {
        "statusCode": status,
        "headers": {**CORS_HEADERS, "Content-Type": "application/json"},
        "body": json.dumps(body, ensure_ascii=False, default=str),
    }


def exc_detail(stage: str, e: Exception, extra: dict = None) -> dict:
    """Возвращает структурированные детали исключения для диагностики."""
    trace_id = uuid.uuid4().hex[:12]
    d = {
        "trace_id": trace_id,
        "stage": stage,
        "error_type": type(e).__name__,
        "message": str(e),
    }
    pgcode = getattr(e, "pgcode", None)
    pgerror = getattr(e, "pgerror", None)
    diag = getattr(e, "diag", None)
    if pgcode:
        d["pgcode"] = pgcode
    if pgerror:
        d["pgerror"] = pgerror.strip() if pgerror else None
    if diag:
        d["pg_primary"] = getattr(diag, "message_primary", None)
        d["pg_schema"] = getattr(diag, "schema_name", None)
        d["pg_table"] = getattr(diag, "table_name", None)
    if isinstance(e, ClientError):
        err = e.response.get("Error", {})
        d["aws_code"] = err.get("Code")
        d["aws_message"] = err.get("Message")
    if extra:
        d.update(extra)
    log.error(json.dumps({**d, "tb": traceback.format_exc()}, ensure_ascii=False, default=str))
    return d


def get_user(conn, session_id: str):
    if not session_id:
        return None
    schema = get_schema()
    cur = conn.cursor()
    cur.execute(
        f"SELECT u.id, u.email FROM {schema}.sessions s "
        f"JOIN {schema}.users u ON u.id = s.user_id "
        f"WHERE s.id = %s AND s.expires_at > NOW()",
        (session_id,),
    )
    row = cur.fetchone()
    return {"id": str(row[0]), "email": row[1]} if row else None


# ------------------------------------------------------------------ #
#  PPTX text extractor                                                #
# ------------------------------------------------------------------ #

def extract_pptx_text(pptx_bytes: bytes) -> list:
    """Возвращает список {"slide": N, "title": str, "text": str}."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(pptx_bytes))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            title = ""
            texts = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                t = shape.text_frame.text.strip()
                if not t:
                    continue
                if shape.shape_type == 13:  # picture
                    continue
                is_title = False
                if not title and len(t) < 120:
                    try:
                        ph = shape.placeholder_format
                        if ph is not None and ph.idx == 0:
                            is_title = True
                    except Exception:
                        pass
                    if not is_title and shape.shape_id in (2, 3):
                        is_title = True
                if is_title:
                    title = t
                else:
                    texts.append(t)
            # Notes
            notes_text = ""
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            full_text = "\n".join(texts)
            slides.append({
                "slide": i,
                "title": title or f"Слайд {i}",
                "text": full_text,
                "notes": notes_text,
            })
        return slides
    except Exception as e:
        return [{"slide": 1, "title": "Ошибка", "text": f"[Не удалось прочитать PPTX: {e}]", "notes": ""}]


# ------------------------------------------------------------------ #
#  YandexGPT call                                                     #
# ------------------------------------------------------------------ #

def call_gpt(messages: list) -> str:
    import urllib.request
    api_key = os.environ.get("YANDEX_GPT_API_KEY", "")
    folder_id = os.environ.get("YANDEX_FOLDER_ID", "")
    if not api_key or not folder_id:
        return "[AI недоступен]"
    payload = json.dumps({
        "modelUri": f"gpt://{folder_id}/yandexgpt/latest",
        "completionOptions": {"stream": False, "temperature": 0.3, "maxTokens": 6000},
        "messages": [{"role": m["role"], "text": m["content"]} for m in messages],
    }).encode()
    req = urllib.request.Request(
        "https://llm.api.cloud.yandex.net/foundationModels/v1/completion",
        data=payload,
        headers={"Authorization": f"Api-Key {api_key}", "Content-Type": "application/json"},
    )
    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            result = json.loads(resp.read())
            return result["result"]["alternatives"][0]["message"]["text"]
    except Exception as e:
        return f"[Ошибка AI: {e}]"


# ------------------------------------------------------------------ #
#  Audit engine                                                       #
# ------------------------------------------------------------------ #

ISSUE_TYPES = [
    "missing_required_topic",
    "missing_required_slide",
    "contradiction_with_source",
    "unsupported_claim",
    "terminology_mismatch",
    "number_mismatch",
    "structure_noncompliance",
    "criteria_noncompliance",
    "template_content_leak",
    "weak_justification",
    "redundancy_or_irrelevance",
]

SYSTEM_PROMPT = """Ты — эксперт по аудиту презентаций. Твоя задача:
1. Проанализировать слайды презентации.
2. Сравнить содержимое с предоставленными документами по их ролям.
3. Найти проблемы: отсутствующие разделы, противоречия, ненадёжные утверждения, несоответствия критериям.
4. Для каждой проблемы указать: slide, issue_type, severity, объяснение, фрагмент из презентации, фрагмент из документа.

Роли документов:
- STANDARD / CRITERIA: нормативные требования — главный приоритет
- SOURCE: источник фактов — проверять на соответствие
- TEMPLATE / EXAMPLE: только для структуры, не источник истины
- MATERIAL: дополнительный контекст

Правила:
- Не галлюцинировать несоответствия. Низкая уверенность → пометить "нужна ручная проверка".
- Каждое замечание — конкретное, с цитатами.
- Не ругать за расхождение с TEMPLATE/EXAMPLE по предметному содержанию.

Отвечай ТОЛЬКО валидным JSON без пояснений и markdown-обёрток.
"""


def run_audit(pptx_slides: list, documents: list) -> dict:
    """Запускает AI-анализ и возвращает structured audit result."""

    # Формируем контекст документов
    doc_ctx_parts = []
    for doc in sorted(documents, key=lambda d: ROLE_PRIORITY.get(d.get("role", "material"), 99)):
        role = doc.get("role", "material").upper()
        name = doc.get("name", "Документ")
        text = (doc.get("text") or "")[:4000]
        instruction = doc.get("instruction") or ""
        part = f"=== [{role}] {name} ===\n{text}"
        if instruction:
            part += f"\n[Инструкция: {instruction}]"
        doc_ctx_parts.append(part)
    doc_context = "\n\n".join(doc_ctx_parts)

    # Контекст слайдов
    slides_ctx = "\n\n".join(
        f"--- Слайд {s['slide']}: {s['title']} ---\n{s['text']}"
        + (f"\n[Notes: {s['notes']}]" if s.get("notes") else "")
        for s in pptx_slides
    )

    # Список документов для явного указания в промпте
    docs_list_txt = "\n".join(
        f"- [{doc.get('role','material').upper()}] {doc.get('name','Документ')}"
        for doc in sorted(documents, key=lambda d: ROLE_PRIORITY.get(d.get("role","material"), 99))
    )

    user_prompt = f"""ДОКУМЕНТЫ ДЛЯ ПРОВЕРКИ ({len(documents)} шт.):
{docs_list_txt}

{doc_context}

СЛАЙДЫ ПРЕЗЕНТАЦИИ ({len(pptx_slides)} шт.):
{slides_ctx}

Проверь презентацию и верни JSON в формате:
{{
  "audit_summary": {{
    "total_slides": <число>,
    "total_issues": <число>,
    "critical_count": <число>,
    "high_count": <число>,
    "medium_count": <число>,
    "low_count": <число>,
    "compliance_score": <0-100>,
    "matched_count": <число критериев выполнено>,
    "partial_count": <число критериев частично>,
    "key_risks": ["..."]
  }},
  "criteria": [
    {{
      "criterion_id": "C001",
      "role": "название роли документа",
      "title": "Краткое название критерия",
      "description": "Полная формулировка требования",
      "source_document": "название файла",
      "source_quote": "точная цитата из документа, на которой основан критерий"
    }}
  ],
  "findings": [
    {{
      "issue_id": "F001",
      "criterion_id": "C001",
      "severity": "critical|high|medium|low",
      "slide_index": <номер>,
      "slide_title": "...",
      "issue_type": "<тип из списка: {', '.join(ISSUE_TYPES)}>",
      "short_title": "...",
      "explanation": "...",
      "what_required": "Что требовал документ",
      "what_found": "Что найдено в презентации",
      "gap_description": "В чём конкретно расхождение",
      "evidence_from_presentation": "цитата из слайда",
      "evidence_from_source_docs": "цитата из документа",
      "related_document_name": "название документа",
      "violated_criterion": "...",
      "suggested_fix": "Конкретная правка текста",
      "rationale": "Почему именно такая правка",
      "confidence": "high|medium|low"
    }}
  ],
  "slide_reports": [
    {{
      "slide_index": <N>,
      "slide_title": "...",
      "status": "ok|needs_attention|critical",
      "issue_count": <число>,
      "summary": "Краткое резюме по слайду"
    }}
  ],
  "compliance_matrix": [
    {{
      "criterion_id": "C001",
      "criterion": "Описание требования",
      "source": "Название документа",
      "status": "met|partially_met|not_met|not_checked",
      "slide_index": <N или null>,
      "comment": "..."
    }}
  ],
  "unverified_items": [
    {{
      "criterion_id": "C001",
      "criterion": "Что не удалось проверить",
      "reason": "insufficient_data|ambiguous_criterion|missing_section|no_relevant_slide",
      "reason_text": "Объяснение почему не удалось проверить"
    }}
  ],
  "suggested_changes": [
    {{
      "slide_index": <N>,
      "slide_title": "...",
      "action": "rewrite|add|remove|replace",
      "current_text": "...",
      "proposed_text": "...",
      "rationale": "..."
    }}
  ],
  "warnings": ["..."]
}}"""

    raw = call_gpt([
        {"role": "system", "content": SYSTEM_PROMPT},
        {"role": "user", "content": user_prompt},
    ])

    # Парсим JSON из ответа
    import re
    json_match = re.search(r'\{.*\}', raw, re.DOTALL)
    if json_match:
        try:
            return json.loads(json_match.group(0))
        except Exception:
            pass

    # Fallback
    return {
        "audit_summary": {
            "total_slides": len(pptx_slides),
            "total_issues": 0,
            "critical_count": 0,
            "high_count": 0,
            "medium_count": 0,
            "low_count": 0,
            "compliance_score": None,
            "matched_count": 0,
            "partial_count": 0,
            "key_risks": [],
        },
        "criteria": [],
        "findings": [],
        "slide_reports": [],
        "compliance_matrix": [],
        "unverified_items": [],
        "suggested_changes": [],
        "warnings": [f"AI вернул нечитаемый ответ: {raw[:200]}"],
    }


# ------------------------------------------------------------------ #
#  Handler                                                            #
# ------------------------------------------------------------------ #

def get_s3():
    import boto3
    return boto3.client(
        "s3",
        endpoint_url="https://bucket.poehali.dev",
        aws_access_key_id=os.environ["AWS_ACCESS_KEY_ID"],
        aws_secret_access_key=os.environ["AWS_SECRET_ACCESS_KEY"],
    )


def handler(event: dict, context) -> dict:
    """Аудит презентаций: проверка PPTX на соответствие документам."""
    if event.get("httpMethod") == "OPTIONS":
        return {"statusCode": 200, "headers": {**CORS_HEADERS, "Content-Type": "application/json"}, "body": "{}"}

    if event.get("httpMethod") != "POST":
        return err_resp("Только POST", 405)

    body = {}
    try:
        body = json.loads(event.get("body") or "{}")
    except Exception:
        return err_resp("Невалидный JSON", 400)

    session_id = event.get("headers", {}).get("X-Session-Id", "")
    conn = get_db()
    try:
        schema = get_schema()
        log.info(f"audit: schema={schema!r}, action={body.get('action')!r}")
        try:
            user = get_user(conn, session_id)
        except Exception as e:
            conn.rollback()
            detail = exc_detail("get_user", e, {"schema": schema})
            return err_resp("Ошибка авторизации", 500, detail=detail)
        if not user:
            return err_resp("Требуется авторизация", 401)

        action = body.get("action", "")

        # ---- audit.setup_cors (одноразовый admin endpoint) ----
        if action == "audit.setup_cors":
            expected_token = os.environ.get("CORS_SETUP_TOKEN", "")
            provided_token = body.get("token", "")
            if not expected_token or provided_token != expected_token:
                return err_resp("Неверный токен", 403)

            s3 = get_s3()
            cors_config = {
                "CORSRules": [{
                    "AllowedOrigins": ["*"],
                    "AllowedMethods": ["PUT", "GET", "HEAD", "DELETE"],
                    "AllowedHeaders": ["*"],
                    "ExposeHeaders": ["ETag"],
                    "MaxAgeSeconds": 3600,
                }]
            }
            s3.put_bucket_cors(Bucket="files", CORSConfiguration=cors_config)
            applied = s3.get_bucket_cors(Bucket="files")
            log.info(f"audit.setup_cors: applied by user={user['email']}")
            return ok_resp({"message": "CORS настроен успешно", "cors": applied.get("CORSRules", [])})

        # ---- audit.upload_init ----
        # Начинает сессию чанковой загрузки, возвращает session_id.
        if action == "audit.upload_init":
            project_id = body.get("project_id")
            filename = body.get("filename", "presentation.pptx")
            total_size = int(body.get("total_size", 0))

            if not project_id:
                return err_resp("Нужен project_id")

            clean_name = (filename or "").lower().strip()
            if not clean_name.endswith(".pptx"):
                return err_resp("Поддерживается только формат PPTX (.pptx)")

            MAX_SIZE = 200 * 1024 * 1024
            if total_size > MAX_SIZE:
                return err_resp(f"Файл слишком большой. Максимум — 200 МБ")

            cur = conn.cursor()
            cur.execute(
                f"SELECT role FROM {schema}.project_members WHERE project_id = %s AND user_id = %s",
                (project_id, user["id"]),
            )
            if not cur.fetchone():
                return err_resp("Нет доступа к проекту", 403)

            session_id = "upl_" + uuid.uuid4().hex
            PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            s3_key = f"audit_uploads/{user['id']}/{project_id}/{session_id}.pptx"

            cur.execute(
                f"""INSERT INTO {schema}.audit_uploads
                    (id, project_id, user_id, filename, content_type, size_bytes_expected, s3_key, status, expires_at)
                    VALUES (%s, %s, %s, %s, %s, %s, %s, 'uploading', NOW() + INTERVAL '1 hour')""",
                (session_id, project_id, user["id"], filename, PPTX_MIME, total_size, s3_key),
            )
            conn.commit()

            log.info(f"audit.upload_init: session_id={session_id}, size={total_size}")
            return ok_resp({"session_id": session_id})

        # ---- audit.upload_chunk ----
        # Принимает один чанк base64, сохраняет как отдельный объект в S3.
        if action == "audit.upload_chunk":
            import base64
            session_id = body.get("session_id")
            chunk_b64 = body.get("chunk_b64", "")
            chunk_index = int(body.get("chunk_index", 0))

            if not session_id or not chunk_b64:
                return err_resp("Нужны session_id и chunk_b64")

            try:
                chunk_bytes = base64.b64decode(chunk_b64)
            except Exception:
                return err_resp("Невалидный base64")

            cur = conn.cursor()
            try:
                cur.execute(
                    f"SELECT s3_key, status FROM {schema}.audit_uploads WHERE id = %s AND user_id = %s",
                    (session_id, user["id"]),
                )
            except Exception as e:
                conn.rollback()
                detail = exc_detail("upload_chunk.select", e, {"schema": schema, "session_id": session_id})
                return err_resp("Ошибка БД при поиске сессии", 500, detail=detail)

            row = cur.fetchone()
            if not row:
                return err_resp(f"Сессия не найдена: {session_id}", 404)
            s3_key, status = row
            if status not in ("uploading",):
                return err_resp(f"Неверный статус сессии: {status}")

            # Сохраняем чанк как отдельный объект: s3_key.chunk.N
            chunk_key = f"{s3_key}.chunk.{chunk_index:04d}"
            try:
                s3 = get_s3()
                s3.put_object(Bucket="files", Key=chunk_key, Body=chunk_bytes)
            except Exception as e:
                detail = exc_detail("upload_chunk.s3_put", e, {"chunk_key": chunk_key, "chunk_index": chunk_index})
                return err_resp("Ошибка S3 при сохранении чанка", 500, detail=detail)

            log.info(f"audit.upload_chunk: session={session_id}, chunk={chunk_index}, key={chunk_key}, size={len(chunk_bytes)}")
            return ok_resp({"chunk_index": chunk_index, "size": len(chunk_bytes)})

        # ---- audit.upload_complete ----
        # Склеивает все чанки из S3 в один файл и регистрирует как готовый.
        if action == "audit.upload_complete":
            import base64
            session_id = body.get("session_id")
            total_chunks = int(body.get("total_chunks", 0))

            if not session_id or not total_chunks:
                return err_resp("Нужны session_id и total_chunks")

            cur = conn.cursor()
            try:
                cur.execute(
                    f"SELECT s3_key FROM {schema}.audit_uploads WHERE id = %s AND user_id = %s",
                    (session_id, user["id"]),
                )
            except Exception as e:
                conn.rollback()
                detail = exc_detail("upload_complete.select", e, {"schema": schema, "session_id": session_id})
                return err_resp("Ошибка БД при поиске сессии", 500, detail=detail)

            row = cur.fetchone()
            if not row:
                return err_resp(f"Сессия не найдена: {session_id}", 404)
            s3_key = row[0]

            s3 = get_s3()

            # Скачиваем и склеиваем все чанки
            file_parts = []
            for i in range(total_chunks):
                chunk_key = f"{s3_key}.chunk.{i:04d}"
                try:
                    obj = s3.get_object(Bucket="files", Key=chunk_key)
                    file_parts.append(obj["Body"].read())
                    s3.delete_object(Bucket="files", Key=chunk_key)
                except Exception as e:
                    detail = exc_detail("upload_complete.get_chunk", e, {
                        "chunk_index": i, "chunk_key": chunk_key,
                        "total_chunks": total_chunks, "chunks_got": len(file_parts),
                    })
                    return err_resp(f"Ошибка S3: не найден чанк {i}", 500, detail=detail)

            file_bytes = b"".join(file_parts)

            PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            try:
                s3.put_object(Bucket="files", Key=s3_key, Body=file_bytes, ContentType=PPTX_MIME)
            except Exception as e:
                detail = exc_detail("upload_complete.put_final", e, {"s3_key": s3_key, "size": len(file_bytes)})
                return err_resp("Ошибка S3 при сохранении файла", 500, detail=detail)

            cur.execute(
                f"UPDATE {schema}.audit_uploads SET status = 'ready', size_bytes_actual = %s WHERE id = %s",
                (len(file_bytes), session_id),
            )
            conn.commit()

            log.info(f"audit.upload_complete: session={session_id}, s3_key={s3_key}, total_size={len(file_bytes)}, chunks={total_chunks}")
            return ok_resp({"upload_id": session_id, "s3_key": s3_key})

        # ---- audit.upload (устарело — оставлено для обратной совместимости) ----
        if action == "audit.upload":
            return err_resp("Используйте audit.upload_init / upload_chunk / upload_complete", 410)

        # ---- audit.prepare_upload (устарело) ----
        if action == "audit.prepare_upload":
            return err_resp("Используйте audit.upload_init для загрузки файлов", 410)

        # ---- audit.run ----
        if action == "audit.run":
            project_id = body.get("project_id")
            upload_id = body.get("upload_id")
            documents = body.get("documents") or []

            if not project_id or not upload_id:
                return err_resp("Нужны project_id и upload_id")

            cur = conn.cursor()
            cur.execute(
                f"SELECT role FROM {schema}.project_members WHERE project_id = %s AND user_id = %s",
                (project_id, user["id"]),
            )
            if not cur.fetchone():
                return err_resp("Нет доступа к проекту", 403)

            # Проверяем upload-сессию
            cur.execute(
                f"""SELECT id, s3_key, filename, size_bytes_expected, status, expires_at
                    FROM {schema}.audit_uploads
                    WHERE id = %s AND user_id = %s AND project_id = %s""",
                (upload_id, user["id"], project_id),
            )
            upl = cur.fetchone()
            if not upl:
                return err_resp("Upload не найден или нет доступа", 404)

            upl_id, pptx_s3_key, source_filename, size_bytes_expected, upl_status, expires_at = upl

            if upl_status == "expired":
                return err_resp("Срок загрузки истёк, загрузите файл заново")
            if upl_status == "consumed":
                return err_resp("Этот upload уже использован")

            # HEAD-валидация: проверяем что объект реально существует и размер совпадает
            s3 = get_s3()
            try:
                head = s3.head_object(Bucket="files", Key=pptx_s3_key)
                actual_size = head["ContentLength"]
                log.info(f"audit.run: head_object ok, upload_id={upload_id}, key={pptx_s3_key}, actual={actual_size}, expected={size_bytes_expected}")
                if actual_size == 0:
                    return err_resp("Файл загружен пустым — попробуйте загрузить снова")
                # Обновляем фактический размер
                cur.execute(
                    f"UPDATE {schema}.audit_uploads SET size_bytes_actual=%s, status='uploaded' WHERE id=%s",
                    (actual_size, upload_id),
                )
            except Exception as e:
                log.error(f"audit.run: head_object failed, upload_id={upload_id}, key={pptx_s3_key}, err={e}")
                return err_resp("Файл не найден в хранилище — браузер ещё не завершил загрузку или файл не был загружен. Попробуйте снова")

            # Скачиваем PPTX
            try:
                obj = s3.get_object(Bucket="files", Key=pptx_s3_key)
                pptx_bytes = obj["Body"].read()
                pptx_size = len(pptx_bytes)
                log.info(f"audit.run: downloaded {pptx_size} bytes, magic={pptx_bytes[:4].hex() if pptx_bytes else 'empty'}")
            except Exception as e:
                log.error(f"audit.run: get_object failed: {e}, key={pptx_s3_key}")
                cur.execute(
                    f"UPDATE {schema}.audit_uploads SET status='failed', error_message=%s WHERE id=%s",
                    (str(e), upload_id),
                )
                conn.commit()
                return err_resp("Не удалось прочитать файл из хранилища — загрузите его заново")

            if pptx_size == 0:
                return err_resp("Файл пустой — загрузите его заново")

            if not pptx_bytes.startswith(b"PK\x03\x04"):
                log.error(f"audit.run: not valid ZIP/PPTX, bytes={pptx_bytes[:16].hex()}")
                return err_resp("Файл загружен, но не является валидным PPTX. Убедитесь что файл не повреждён")

            # Читаем документы из БД по ролям из запроса
            role_map = {str(d["document_id"]): d.get("role", "material") for d in documents if isinstance(d, dict) and d.get("document_id")}
            instr_map = {str(d["document_id"]): d.get("instruction", "") for d in documents if isinstance(d, dict) and d.get("document_id")}

            cur.execute(
                f"""SELECT id, original_name, file_type, extracted_text,
                           COALESCE(default_ai_role, 'material') as default_role
                    FROM {schema}.documents
                    WHERE project_id = %s AND archived_at IS NULL
                    ORDER BY created_at""",
                (project_id,),
            )
            db_docs = cur.fetchall()
            doc_list = []
            for r in db_docs:
                doc_id, orig_name, ftype, text, def_role = r
                if not text:
                    continue
                doc_list.append({
                    "name": orig_name,
                    "role": role_map.get(str(doc_id), def_role),
                    "text": text,
                    "instruction": instr_map.get(str(doc_id), ""),
                })

            # Извлекаем текст слайдов
            pptx_slides = extract_pptx_text(pptx_bytes)
            log.info(f"audit.run: extracted {len(pptx_slides)} slides")
            if not pptx_slides or (len(pptx_slides) == 1 and "Ошибка" in pptx_slides[0].get("title", "")):
                log.error(f"audit.run: extract failed: {pptx_slides}")
                return err_resp("Не удалось прочитать презентацию. Убедитесь, что файл не повреждён и является корректным .pptx")

            # Помечаем upload как consumed
            cur.execute(
                f"UPDATE {schema}.audit_uploads SET status='consumed', consumed_at=NOW(), size_bytes_actual=%s WHERE id=%s",
                (pptx_size, upload_id),
            )

            # Запускаем AI-анализ
            audit_result = run_audit(pptx_slides, doc_list)
            audit_result["slide_count"] = len(pptx_slides)
            audit_result["document_count"] = len(doc_list)
            # Сохраняем точный список документов этого аудита (имя + роль)
            audit_result["documents_used"] = [
                {"name": doc["name"], "role": doc["role"]}
                for doc in doc_list
            ]

            # Сохраняем в БД
            cur.execute(
                f"""INSERT INTO {schema}.audit_runs
                    (project_id, user_id, slide_count, doc_count, result_json, status,
                     source_pptx_s3_key, source_filename, source_size_bytes)
                    VALUES (%s, %s, %s, %s, %s, 'done', %s, %s, %s) RETURNING id""",
                (
                    project_id, user["id"],
                    len(pptx_slides), len(doc_list),
                    json.dumps(audit_result, ensure_ascii=False, default=str),
                    pptx_s3_key, source_filename, pptx_size,
                ),
            )
            audit_id = cur.fetchone()[0]
            conn.commit()

            return ok_resp({"audit_id": audit_id, "result": audit_result})

        # ---- audit.get ----
        if action == "audit.get":
            audit_id = body.get("audit_id")
            if not audit_id:
                return err_resp("Нужен audit_id")
            cur = conn.cursor()
            cur.execute(
                f"""SELECT ar.id, ar.project_id, ar.slide_count, ar.doc_count,
                           ar.result_json, ar.status, ar.created_at
                    FROM {schema}.audit_runs ar
                    WHERE ar.id = %s AND ar.user_id = %s""",
                (int(audit_id), user["id"]),
            )
            row = cur.fetchone()
            if not row:
                return err_resp("Аудит не найден", 404)
            result = json.loads(row[4]) if row[4] else {}
            return ok_resp({
                "audit_id": row[0], "project_id": row[1],
                "slide_count": row[2], "doc_count": row[3],
                "result": result, "status": row[5],
                "created_at": str(row[6]),
            })

        # ---- audit.list ----
        if action == "audit.list":
            project_id = body.get("project_id")
            if not project_id:
                return err_resp("Нужен project_id")
            cur = conn.cursor()
            cur.execute(
                f"SELECT role FROM {schema}.project_members WHERE project_id = %s AND user_id = %s",
                (project_id, user["id"]),
            )
            if not cur.fetchone():
                return err_resp("Нет доступа", 403)
            cur.execute(
                f"""SELECT id, slide_count, doc_count, status, created_at,
                           (result_json::json->'audit_summary'->>'compliance_score') as score,
                           (result_json::json->'audit_summary'->>'total_issues') as issues
                    FROM {schema}.audit_runs
                    WHERE project_id = %s AND user_id = %s
                    ORDER BY created_at DESC LIMIT 20""",
                (project_id, user["id"]),
            )
            rows = cur.fetchall()
            return ok_resp({"audits": [
                {
                    "audit_id": r[0], "slide_count": r[1], "doc_count": r[2],
                    "status": r[3], "created_at": str(r[4]),
                    "compliance_score": r[5], "total_issues": r[6],
                }
                for r in rows
            ]})

        # ================================================================
        # audit.build_revision_plan
        # ================================================================
        if action == "audit.build_revision_plan":
            audit_id = int(body.get("audit_id", 0))
            options = body.get("options") or {}   # revision_mode, filters, etc.

            cur = conn.cursor()
            cur.execute(
                f"SELECT result_json, project_id FROM {schema}.audit_runs WHERE id = %s AND user_id = %s",
                (audit_id, user["id"]),
            )
            row = cur.fetchone()
            if not row:
                return err_resp("Аудит не найден", 404)

            audit_result = json.loads(row[0]) if row[0] else {}
            findings = audit_result.get("findings") or []
            suggested = audit_result.get("suggested_changes") or []
            compliance = audit_result.get("compliance_matrix") or []

            # --- Фильтрация findings по опциям ---
            severity_filter = options.get("severity_filter", ["critical", "high"])  # по умолчанию только critical+high
            exclude_low_confidence = options.get("exclude_low_confidence", True)

            applicable = []
            skipped = []
            for f in findings:
                sev = f.get("severity", "low")
                conf = f.get("confidence", "high")
                if sev not in severity_filter:
                    skipped.append({"issue_id": f.get("issue_id"), "reason": f"severity={sev} не в фильтре"})
                    continue
                if exclude_low_confidence and conf == "low":
                    skipped.append({"issue_id": f.get("issue_id"), "reason": "низкая уверенность AI"})
                    continue
                applicable.append(f)

            # --- AI строит revision plan ---
            revision_mode = options.get("revision_mode", "fix_text")  # fix_text | fix_and_add | full_revision
            keep_slide_count = options.get("keep_slide_count", True)
            allow_add_slides = options.get("allow_add_slides", False)
            keep_visuals = options.get("keep_visuals", True)

            findings_txt = json.dumps(applicable, ensure_ascii=False, default=str)
            suggested_txt = json.dumps(suggested, ensure_ascii=False, default=str)
            compliance_txt = json.dumps([c for c in compliance if c.get("status") != "met"], ensure_ascii=False, default=str)

            plan_prompt = f"""Ты — редактор презентаций. На основе результатов аудита составь план исправлений.

ПРИМЕНЯЕМЫЕ ЗАМЕЧАНИЯ ({len(applicable)} шт.):
{findings_txt[:3000]}

ПРЕДЛОЖЕННЫЕ ПРАВКИ:
{suggested_txt[:2000]}

НЕВЫПОЛНЕННЫЕ КРИТЕРИИ:
{compliance_txt[:1000]}

РЕЖИМ ИСПРАВЛЕНИЯ: {revision_mode}
СОХРАНЯТЬ ЧИСЛО СЛАЙДОВ: {keep_slide_count}
РАЗРЕШИТЬ ДОБАВЛЕНИЕ СЛАЙДОВ: {allow_add_slides}
СОХРАНЯТЬ ВИЗУАЛЫ: {keep_visuals}

Составь детальный план исправлений. Верни ТОЛЬКО JSON:
{{
  "revision_plan": [
    {{
      "plan_item_id": "P001",
      "slide_index": <N>,
      "slide_title": "...",
      "change_type": "rewrite_text|add_missing_point|remove_unsupported_claim|replace_terminology|add_missing_slide|restructure_slide|update_numbers|mark_for_manual_review",
      "based_on_finding_ids": ["F001", "F002"],
      "problem_summary": "...",
      "proposed_change": "Конкретное изменение текста или структуры",
      "rationale": "Почему это исправление нужно",
      "confidence": "high|medium|low",
      "will_affect_visual": false,
      "visual_action": "keep|needs_review|needs_regeneration|preserve_user_override",
      "requires_user_review": false,
      "priority": 1
    }}
  ],
  "revision_summary": {{
    "total_changes": <N>,
    "slides_affected": [1, 2, 3],
    "will_add_slides": false,
    "expected_improvement": "Краткое описание ожидаемого результата",
    "manual_review_required": ["P003", "P007"]
  }},
  "generate_instruction": "Инструкция для AI генерации исправленной версии в 2-3 абзаца"
}}"""

            plan_raw = call_gpt([
                {"role": "system", "content": "Ты — редактор презентаций. Отвечай ТОЛЬКО валидным JSON без markdown."},
                {"role": "user", "content": plan_prompt},
            ])

            import re as _re
            plan_data = {}
            jm = _re.search(r'\{.*\}', plan_raw, _re.DOTALL)
            if jm:
                try:
                    plan_data = json.loads(jm.group(0))
                except Exception:
                    pass

            if not plan_data:
                plan_data = {
                    "revision_plan": [],
                    "revision_summary": {"total_changes": 0, "expected_improvement": "AI не смог построить план"},
                    "generate_instruction": "",
                }

            plan_data["applicable_findings"] = applicable
            plan_data["skipped_findings"] = skipped
            plan_data["options"] = options

            # Сохраняем план в audit_runs
            cur.execute(
                f"UPDATE {schema}.audit_runs SET revision_plan_json = %s, revision_status = 'plan_ready' WHERE id = %s",
                (json.dumps(plan_data, ensure_ascii=False, default=str), audit_id),
            )
            conn.commit()

            return ok_resp({"audit_id": audit_id, "revision_plan": plan_data})

        # ================================================================
        # audit.create_revision_run — создаёт задание + run на исправление
        # ================================================================
        if action == "audit.create_revision_run":
            audit_id = int(body.get("audit_id", 0))
            task_id = body.get("task_id")
            documents = body.get("documents") or []
            confirmed_plan_items = body.get("confirmed_plan_items")

            cur = conn.cursor()
            cur.execute(
                f"""SELECT result_json, revision_plan_json, project_id, source_pptx_s3_key
                    FROM {schema}.audit_runs WHERE id = %s AND user_id = %s""",
                (audit_id, user["id"]),
            )
            row = cur.fetchone()
            if not row:
                return err_resp("Аудит не найден", 404)

            audit_result = json.loads(row[0]) if row[0] else {}
            plan_data = json.loads(row[1]) if row[1] else {}
            project_id = row[2]
            source_pptx_s3_key = row[3]

            # Берём документы из БД по document_id+role из запроса
            role_map_cr = {str(d["document_id"]): d.get("role", "material") for d in documents if isinstance(d, dict) and d.get("document_id")}
            instr_map_cr = {str(d["document_id"]): d.get("instruction", "") for d in documents if isinstance(d, dict) and d.get("document_id")}
            cur.execute(
                f"""SELECT id, original_name, COALESCE(default_ai_role, 'material'), extracted_text
                    FROM {schema}.documents
                    WHERE project_id = %s AND archived_at IS NULL AND extracted_text IS NOT NULL""",
                (project_id,),
            )
            documents = [
                {"name": r[1], "role": role_map_cr.get(str(r[0]), r[2]), "text": r[3], "instruction": instr_map_cr.get(str(r[0]), "")}
                for r in cur.fetchall()
            ]

            if not plan_data:
                return err_resp("Сначала постройте план исправлений (audit.build_revision_plan)")

            revision_plan = plan_data.get("revision_plan") or []
            if confirmed_plan_items is not None:
                revision_plan = [p for p in revision_plan if p.get("plan_item_id") in confirmed_plan_items]

            # --- Формируем instruction для generate ---
            base_instruction = plan_data.get("generate_instruction", "")
            changes_txt = "\n".join(
                f"- Слайд {p['slide_index']} ({p['change_type']}): {p['proposed_change']}"
                for p in revision_plan[:20]
            )
            findings_for_gen = json.dumps(
                audit_result.get("findings", [])[:10], ensure_ascii=False, default=str
            )

            final_instruction = f"""ЭТО ИСПРАВЛЕННАЯ ВЕРСИЯ НА ОСНОВЕ АУДИТА.

Исходный аудит выявил следующие проблемы (применяются):
{findings_for_gen[:2000]}

КОНКРЕТНЫЙ ПЛАН ИСПРАВЛЕНИЙ ПО СЛАЙДАМ:
{changes_txt}

{base_instruction}

ПРАВИЛА:
1. Исправляй ТОЛЬКО то, что указано в плане — не переписывай всё заново.
2. Сохраняй общую структуру презентации, если не указано иное.
3. При исправлении терминологии — используй термины из документов-источников.
4. Не добавляй новые утверждения без основания в документах.
5. Нумеруй слайды явно: "Слайд 1:", "Слайд 2:", и т.д."""

            # Если есть исходное задание — создаём ревизию
            if task_id:
                # Создаём новый run как ревизию исходного задания
                cur.execute(
                    f"""SELECT version_number FROM {schema}.generation_runs
                        WHERE task_id = %s ORDER BY version_number DESC LIMIT 1""",
                    (int(task_id),),
                )
                vr = cur.fetchone()
                next_version = (vr[0] + 1) if vr else 1

                cur.execute(
                    f"""INSERT INTO {schema}.generation_runs
                        (task_id, created_by, version_number, input_prompt, system_constraints, status)
                        VALUES (%s, %s, %s, %s, %s, 'running') RETURNING id""",
                    (int(task_id), user["id"], next_version, final_instruction[:2000], "revision_from_audit", ),
                )
                new_run_id = cur.fetchone()[0]

                # Вызываем AI для генерации исправленной версии
                doc_ctx = "\n\n".join(
                    f"=== [{d.get('role','material').upper()}] {d.get('name','Документ')} ===\n{(d.get('text') or '')[:3000]}"
                    for d in sorted(documents, key=lambda d: ROLE_PRIORITY.get(d.get("role","material"), 99))
                )

                ai_content = call_gpt([
                    {"role": "system", "content": f"Ты — профессиональный редактор презентаций.\n\nДОКУМЕНТЫ:\n{doc_ctx[:4000]}"},
                    {"role": "user", "content": final_instruction},
                ])

                # Строим revision_meta для result_json
                applied_ids = [p.get("plan_item_id") for p in revision_plan]
                all_ids = [p.get("plan_item_id") for p in (plan_data.get("revision_plan") or [])]
                skipped_ids = [i for i in all_ids if i not in applied_ids]

                # Визуальные изменения
                visual_changes = []
                for p in revision_plan:
                    if p.get("will_affect_visual") or p.get("visual_action") not in (None, "keep"):
                        visual_changes.append({
                            "slide_index": p["slide_index"],
                            "visual_action": p.get("visual_action", "keep"),
                            "reason": p.get("problem_summary", ""),
                        })

                revision_meta = {
                    "source_audit_run_id": audit_id,
                    "revision_mode": plan_data.get("options", {}).get("revision_mode", "fix_text"),
                    "applied_finding_ids": [
                        f["issue_id"] for f in plan_data.get("applicable_findings", [])
                    ],
                    "applied_plan_item_ids": applied_ids,
                    "skipped_plan_item_ids": skipped_ids,
                    "visual_changes": visual_changes,
                    "revision_plan": revision_plan,
                    "warnings": [],
                }

                result_payload = {
                    "content": ai_content,
                    "version": next_version,
                    "revision_meta": revision_meta,
                }

                cur.execute(
                    f"""UPDATE {schema}.generation_runs
                        SET result_json = %s, output_summary = %s, status = 'done'
                        WHERE id = %s""",
                    (json.dumps(result_payload, ensure_ascii=False, default=str),
                     ai_content[:300],
                     new_run_id),
                )

                # Линкуем audit к новому run
                cur.execute(
                    f"""UPDATE {schema}.audit_runs
                        SET revision_run_id = %s, revision_status = 'revision_done'
                        WHERE id = %s""",
                    (new_run_id, audit_id),
                )
                conn.commit()

                return ok_resp({
                    "audit_id": audit_id,
                    "run_id": new_run_id,
                    "task_id": int(task_id),
                    "version": next_version,
                    "content": ai_content,
                    "revision_meta": revision_meta,
                })
            else:
                # Нет исходного задания — возвращаем готовый текст без run
                doc_ctx = "\n\n".join(
                    f"=== [{d.get('role','material').upper()}] {d.get('name','Документ')} ===\n{(d.get('text') or '')[:3000]}"
                    for d in sorted(documents, key=lambda d: ROLE_PRIORITY.get(d.get("role","material"), 99))
                )
                ai_content = call_gpt([
                    {"role": "system", "content": f"Ты — профессиональный редактор презентаций.\n\nДОКУМЕНТЫ:\n{doc_ctx[:4000]}"},
                    {"role": "user", "content": final_instruction},
                ])

                cur.execute(
                    f"UPDATE {schema}.audit_runs SET revision_status = 'revision_done' WHERE id = %s",
                    (audit_id,),
                )
                conn.commit()

                applied_ids_no_task = [p.get("plan_item_id") for p in revision_plan]
                all_ids_no_task = [p.get("plan_item_id") for p in (plan_data.get("revision_plan") or [])]
                skipped_ids_no_task = [i for i in all_ids_no_task if i not in applied_ids_no_task]
                return ok_resp({
                    "audit_id": audit_id,
                    "content": ai_content,
                    "revision_meta": {
                        "source_audit_run_id": audit_id,
                        "applied_plan_item_ids": applied_ids_no_task,
                        "skipped_plan_item_ids": skipped_ids_no_task,
                        "applied_finding_ids": [f["issue_id"] for f in plan_data.get("applicable_findings", [])],
                        "visual_changes": [],
                        "warnings": [],
                    },
                })

        # ================================================================
        # audit.download_revised — скачать исправленный PPTX
        # Применяет правки из revision_plan к оригинальному файлу
        # ================================================================
        if action == "audit.download_revised":
            audit_id = int(body.get("audit_id", 0))

            cur = conn.cursor()
            cur.execute(
                f"""SELECT source_pptx_s3_key, revision_plan_json, revision_run_id
                    FROM {schema}.audit_runs WHERE id = %s AND user_id = %s""",
                (audit_id, user["id"]),
            )
            row = cur.fetchone()
            if not row:
                return err_resp("Аудит не найден", 404)

            source_key, plan_json_raw, revision_run_id = row
            if not source_key:
                return err_resp("Исходный файл не найден")

            plan_data = json.loads(plan_json_raw) if plan_json_raw else {}
            revision_plan = plan_data.get("revision_plan") or []

            # Загружаем оригинальный PPTX
            try:
                s3 = get_s3()
                obj = s3.get_object(Bucket="files", Key=source_key)
                pptx_bytes = obj["Body"].read()
            except Exception as e:
                detail = exc_detail("download_revised.get_s3", e, {"key": source_key})
                return err_resp("Не удалось загрузить исходный файл", 500, detail=detail)

            # Применяем правки из плана
            try:
                from pptx import Presentation
                from pptx.util import Pt
                import copy

                prs = Presentation(io.BytesIO(pptx_bytes))
                applied = 0

                for item in revision_plan:
                    slide_index = item.get("slide_index")
                    change_type = item.get("change_type", "")
                    proposed = item.get("proposed_change", "")
                    target_text = item.get("target_text", "")

                    if not proposed or slide_index is None:
                        continue

                    # Индекс слайда 1-based
                    idx = int(slide_index) - 1
                    if idx < 0 or idx >= len(prs.slides):
                        continue

                    slide = prs.slides[idx]

                    if change_type in ("text_replace", "text_add", "fix"):
                        # Ищем фрейм с целевым текстом и заменяем
                        replaced = False
                        for shape in slide.shapes:
                            if not shape.has_text_frame:
                                continue
                            for para in shape.text_frame.paragraphs:
                                full = "".join(r.text for r in para.runs)
                                if target_text and target_text[:40].lower() in full.lower():
                                    # Заменяем текст первого run, остальные очищаем
                                    if para.runs:
                                        para.runs[0].text = proposed
                                        for r in para.runs[1:]:
                                            r.text = ""
                                    replaced = True
                                    applied += 1
                                    break
                            if replaced:
                                break

                        # Если target не найден — ищем заголовок слайда и добавляем после
                        if not replaced and change_type == "text_add":
                            for shape in slide.shapes:
                                if shape.has_text_frame and shape.shape_id in (2, 3):
                                    tf = shape.text_frame
                                    p = tf.add_paragraph()
                                    p.text = proposed
                                    applied += 1
                                    break

                    elif change_type == "title_replace":
                        for shape in slide.shapes:
                            try:
                                ph = shape.placeholder_format
                                if ph is not None and ph.idx == 0 and shape.has_text_frame:
                                    if shape.text_frame.paragraphs:
                                        shape.text_frame.paragraphs[0].runs[0].text = proposed if shape.text_frame.paragraphs[0].runs else ""
                                    applied += 1
                                    break
                            except Exception:
                                pass

                # Сохраняем исправленный PPTX
                buf = io.BytesIO()
                prs.save(buf)
                revised_bytes = buf.getvalue()

            except Exception as e:
                detail = exc_detail("download_revised.apply_fixes", e)
                return err_resp("Ошибка при применении правок", 500, detail=detail)

            # Сохраняем в S3 и отдаём presigned URL (CDN недоступен, используем прямой S3)
            revised_key = source_key.replace(".pptx", f"_revised_{audit_id}.pptx")
            try:
                PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                s3.put_object(Bucket="files", Key=revised_key, Body=revised_bytes, ContentType=PPTX_MIME)
                # Presigned URL — работает напрямую, не зависит от CDN
                download_url = s3.generate_presigned_url(
                    "get_object",
                    Params={"Bucket": "files", "Key": revised_key},
                    ExpiresIn=3600,
                )
            except Exception as e:
                detail = exc_detail("download_revised.put_s3", e)
                return err_resp("Ошибка сохранения исправленного файла", 500, detail=detail)

            log.info(f"audit.download_revised: audit_id={audit_id}, applied={applied}, size={len(revised_bytes)}, key={revised_key}")

            return ok_resp({
                "audit_id": audit_id,
                "applied_count": applied,
                "filename": source_key.split("/")[-1].replace(".pptx", "_revised.pptx"),
                "cdn_url": download_url,
            })

        # ================================================================
        # audit.get_revision_status
        # ================================================================
        if action == "audit.get_revision_status":
            audit_id = int(body.get("audit_id", 0))
            cur = conn.cursor()
            cur.execute(
                f"""SELECT revision_status, revision_plan_json, revision_run_id, reaudit_result_json
                    FROM {schema}.audit_runs WHERE id = %s AND user_id = %s""",
                (audit_id, user["id"]),
            )
            row = cur.fetchone()
            if not row:
                return err_resp("Аудит не найден", 404)

            plan = json.loads(row[1]) if row[1] else None
            reaudit = json.loads(row[3]) if row[3] else None

            return ok_resp({
                "audit_id": audit_id,
                "revision_status": row[0],
                "revision_plan": plan,
                "revision_run_id": row[2],
                "reaudit_result": reaudit,
            })

        # ================================================================
        # audit.run_reaudit — повторный аудит после ревизии
        # Берёт PPTX из source_pptx_s3_key сохранённого в audit_run (не требует UI-state)
        # ================================================================
        if action == "audit.run_reaudit":
            audit_id = int(body.get("audit_id", 0))
            documents = body.get("documents") or []

            cur = conn.cursor()
            cur.execute(
                f"""SELECT project_id, result_json, source_pptx_s3_key
                    FROM {schema}.audit_runs WHERE id = %s AND user_id = %s""",
                (audit_id, user["id"]),
            )
            row = cur.fetchone()
            if not row:
                return err_resp("Аудит не найден", 404)

            project_id_stored, result_json_raw, source_pptx_s3_key = row[0], row[1], row[2]
            original_result = json.loads(result_json_raw) if result_json_raw else {}
            original_score = (original_result.get("audit_summary") or {}).get("compliance_score")

            if not source_pptx_s3_key:
                return err_resp("Исходный файл презентации не найден. Загрузите файл заново для повторной проверки.")

            try:
                s3 = get_s3()
                obj = s3.get_object(Bucket="files", Key=source_pptx_s3_key)
                pptx_bytes = obj["Body"].read()
            except Exception as e:
                return err_resp(f"Не удалось загрузить файл презентации: {e}")

            # Берём документы из БД по document_id+role из запроса
            role_map_r = {str(d["document_id"]): d.get("role", "material") for d in documents if isinstance(d, dict) and d.get("document_id")}
            instr_map_r = {str(d["document_id"]): d.get("instruction", "") for d in documents if isinstance(d, dict) and d.get("document_id")}
            cur.execute(
                f"""SELECT id, original_name, COALESCE(default_ai_role, 'material'), extracted_text
                    FROM {schema}.documents
                    WHERE project_id = %s AND archived_at IS NULL AND extracted_text IS NOT NULL""",
                (project_id_stored,),
            )
            documents = [
                {"name": r[1], "role": role_map_r.get(str(r[0]), r[2]), "text": r[3], "instruction": instr_map_r.get(str(r[0]), "")}
                for r in cur.fetchall()
            ]

            pptx_slides = extract_pptx_text(pptx_bytes)
            new_result = run_audit(pptx_slides, documents)
            new_score = (new_result.get("audit_summary") or {}).get("compliance_score")

            # Дельта
            delta = None
            if original_score is not None and new_score is not None:
                try:
                    delta = int(new_score) - int(original_score)
                except Exception:
                    pass

            reaudit_payload = {
                "audit_result": new_result,
                "score_before": original_score,
                "score_after": new_score,
                "score_delta": delta,
                "issues_before": (original_result.get("audit_summary") or {}).get("total_issues"),
                "issues_after": (new_result.get("audit_summary") or {}).get("total_issues"),
            }

            cur.execute(
                f"""UPDATE {schema}.audit_runs
                    SET reaudit_result_json = %s, revision_status = 'reaudited'
                    WHERE id = %s""",
                (json.dumps(reaudit_payload, ensure_ascii=False, default=str), audit_id),
            )
            conn.commit()

            return ok_resp({"audit_id": audit_id, "reaudit": reaudit_payload})

        return err_resp("Неизвестное действие")

    except Exception as e:
        try:
            conn.rollback()
        except Exception:
            pass
        action_name = body.get("action", "unknown") if isinstance(body, dict) else "unknown"
        detail = exc_detail(f"handler.{action_name}", e)
        return err_resp(f"Ошибка сервера: {detail['message']}", 500, detail=detail)
    finally:
        conn.close()